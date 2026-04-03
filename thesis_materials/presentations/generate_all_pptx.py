"""
Generate all thesis PowerPoint presentations from markdown content.
Run: .venv\Scripts\python.exe thesis_materials/presentations/generate_all_pptx.py

Creates 6 .pptx files in thesis_materials/presentations/
"""

import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── ASU-inspired color scheme ─────────────────────────────────
MAROON   = RGBColor(0x8C, 0x1D, 0x40)  # ASU maroon
GOLD     = RGBColor(0xFF, 0xC6, 0x27)  # ASU gold
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BLACK    = RGBColor(0x19, 0x19, 0x19)
DARK     = RGBColor(0x2D, 0x2D, 0x2D)
GRAY     = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
BLUE     = RGBColor(0x21, 0x71, 0xB5)
GREEN    = RGBColor(0x31, 0xA3, 0x54)
RED      = RGBColor(0xDE, 0x2D, 0x26)
ORANGE   = RGBColor(0xE6, 0x55, 0x0D)

OUT_DIR = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.join(os.path.dirname(OUT_DIR), 'figures', 'generated')


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Send to back
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Support bold prefix with **text**
        if '**' in item:
            parts = item.split('**')
            for j, part in enumerate(parts):
                if not part:
                    continue
                run = p.add_run()
                run.text = part
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.name = 'Calibri'
                run.font.bold = (j % 2 == 1)
        else:
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
    return txBox


def add_table(slide, left, top, width, rows_data, col_widths=None, font_size=12):
    """rows_data: list of lists. First row = header."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                          width, Inches(0.4 * n_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = 'Calibri'
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK
                    if c > 0:
                        paragraph.alignment = PP_ALIGN.CENTER

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = MAROON
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    return table_shape


def make_title_slide(prs, title, subtitle, date="Spring 2026"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, MAROON)

    # Gold accent bar
    add_shape_bg(slide, Inches(0), Inches(3.2), Inches(10), Inches(0.08), GOLD)

    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(8.4), Inches(1.5),
                 title, font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(2.3), Inches(8.4), Inches(0.5),
                 subtitle, font_size=16, color=GOLD, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(8.4), Inches(0.5),
                 "Anudeep Gottapu  |  Arizona State University  |  " + date,
                 font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)
    return slide


def make_section_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, MAROON)
    add_shape_bg(slide, Inches(0), Inches(2.8), Inches(10), Inches(0.06), GOLD)
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(8.4), Inches(0.8),
                 title, font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(3.1), Inches(8.4), Inches(0.5),
                     subtitle, font_size=16, color=GOLD, alignment=PP_ALIGN.CENTER)
    return slide


def make_content_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Title bar
    add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(0.9), MAROON)
    add_text_box(slide, Inches(0.5), Inches(0.1), Inches(9), Inches(0.7),
                 title, font_size=24, bold=True, color=WHITE)

    # Gold accent line
    add_shape_bg(slide, Inches(0), Inches(0.9), Inches(10), Inches(0.04), GOLD)
    return slide


def try_add_image(slide, img_name, left, top, width=None, height=None):
    path = os.path.join(FIG_DIR, img_name)
    if os.path.exists(path):
        kwargs = {'left': left, 'top': top}
        if width:
            kwargs['width'] = width
        if height:
            kwargs['height'] = height
        slide.shapes.add_picture(path, **kwargs)
        return True
    return False


# =================================================================
# WEEK 1: Environment & Demos
# =================================================================
def build_week1():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    make_title_slide(prs,
        "Week 1: Environment Design &\nDemonstration Collection",
        "TwoBlockPick for Legible Robot Motion")

    # Slide 2: Motivation
    s = make_content_slide(prs, "Problem Motivation")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5), [
        "In shared workspaces, robots must communicate intent through motion",
        "A human watching a robot reach for one of two objects needs to predict the goal **early**",
        "Ambiguous trajectories \u2192 human waits \u2192 slower collaboration",
        "Legible trajectories \u2192 early prediction \u2192 efficient teamwork",
        "",
        "**Research Question:** Can we train a generative policy to produce task-successful AND legible trajectories?",
    ], font_size=16)

    # Slide 3: Environment
    s = make_content_slide(prs, "TwoBlockPick Environment")
    add_table(s, Inches(0.5), Inches(1.2), Inches(5.5), [
        ["Parameter", "Value"],
        ["Robot", "Franka Panda 7-DOF"],
        ["Physics", "PyBullet, dt=1/240s"],
        ["Objects", "Two red blocks, 4cm each"],
        ["Left block", "(0.50, +0.07, 0.42)"],
        ["Right block", "(0.50, -0.07, 0.42)"],
        ["Jitter", "\u00b10.015m per block"],
        ["Success", "Lift block above z=0.52m"],
    ], font_size=11)
    add_bullet_list(s, Inches(0.5), Inches(4.3), Inches(9), Inches(1), [
        "**Obs (22-dim):** ee_pos + ee_quat + grip + block poses  |  **Act (5-dim):** \u0394x, \u0394y, \u0394z, \u0394yaw, grip"
    ], font_size=13)

    # Slide 4: Camera
    s = make_content_slide(prs, "Camera Setup (Critical for VLM)")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5), [
        "Camera: yaw=135\u00b0, pitch=-30\u00b0, distance=0.9m, FOV=60\u00b0",
        "**Image-LEFT = World-RIGHT** (and vice versa)",
        "Goal A = image-left = world-right block",
        "Goal B = image-right = world-left block",
        "",
        "This mapping was initially incorrect and caused evaluation errors.",
        "Getting this right was essential for VLM scoring accuracy.",
    ], font_size=16)

    # Slide 5: Demo Strategy
    s = make_content_slide(prs, "Demonstration Collection Strategy")
    add_text_box(s, Inches(0.5), Inches(1.1), Inches(9), Inches(0.4),
                 "400 Expert Demos via Analytical B\u00e9zier Curves", font_size=16, bold=True)
    add_table(s, Inches(0.5), Inches(1.6), Inches(9), [
        ["Style", "Count", "Design"],
        ["Legible (50%)", "200", "Quadratic B\u00e9zier, control point toward goal. Commitment by t\u22480.3"],
        ["Neutral (25%)", "100", "Quadratic B\u00e9zier, control point at y=0. No lateral signal"],
        ["Deceptive (25%)", "100", "Cubic B\u00e9zier, P1 feints wrong side, P2 commits correct"],
    ], font_size=11)
    add_bullet_list(s, Inches(0.5), Inches(3.6), Inches(9), Inches(1.5), [
        "10 block configurations \u00d7 40 demos each",
        "20 pick-left + 20 pick-right per configuration",
        "Each demo: full episode (up to 400 steps) of obs-action pairs",
    ], font_size=14)

    # Slide 6: Data Format
    s = make_content_slide(prs, "Data Format: demos_combined.npz")
    add_table(s, Inches(0.5), Inches(1.2), Inches(7), [
        ["Array", "Shape", "Type"],
        ["obs", "(400, 400, 22)", "float32"],
        ["actions", "(400, 400, 5)", "float32"],
        ["labels", "(400,)", "'left' or 'right'"],
        ["style_labels", "(400,)", "0=legible, 1=neutral, 2=deceptive"],
        ["config_id", "(400,)", "0-9"],
    ], font_size=12)

    # Slide 7: Why Three Styles
    s = make_content_slide(prs, "Why Three Trajectory Styles?")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5), [
        "**Legible demos** are the target behavior we want the policy to produce",
        "**Neutral demos** represent baseline \u201cshortest path\u201d behavior",
        "**Deceptive demos** are worst case \u2014 feinting before committing",
        "Mixed training data lets us evaluate whether the policy learns meaningful style differences",
        "Enables VLM evaluation: does the VLM correctly rank styles by legibility?",
    ], font_size=16)

    # Slide 8: Design Decisions
    s = make_content_slide(prs, "Key Design Decisions")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5), [
        "**Simulation over real robot** \u2014 Enables 400 demos + reproducible evaluation",
        "**Analytical trajectories** \u2014 B\u00e9zier curves give precise control over shape",
        "**Symmetric block placement** \u2014 Legibility comes purely from trajectory shape",
        "**3 trajectory styles** \u2014 Ground truth for evaluating VLM discrimination",
    ], font_size=16)

    path = os.path.join(OUT_DIR, 'week1_environment_and_demos.pptx')
    prs.save(path)
    print(f"  \u2713 {os.path.basename(path)}")


# =================================================================
# WEEK 2: Diffusion Policy
# =================================================================
def build_week2():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    make_title_slide(prs,
        "Week 2: Diffusion Policy\nArchitecture & Training",
        "DDPM-based Action Chunk Prediction")

    # Why Diffusion
    s = make_content_slide(prs, "Why Diffusion Policy?")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5), [
        "**Multimodal distribution modeling** \u2014 Captures legible/neutral/deceptive in one model",
        "**Training-free controllability** \u2014 Iterative denoising allows gradient-based steering",
        "**Chunk-based prediction** \u2014 32-step chunks provide temporal coherence",
        "**State-of-the-art** \u2014 Chi et al. (2023): DDPM policies outperform explicit classes",
        "",
        "**Key insight:** By learning a distribution over trajectories, we can later steer toward legibility.",
    ], font_size=16)

    # DDPM
    s = make_content_slide(prs, "DDPM Training Formulation")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5), [
        "**Forward:** q(x_t | x_0) = N(x_t; \u221a\u03b1\u0305_t x_0, (1-\u03b1\u0305_t)I)",
        "**Objective:** L = E[\u2016\u03b5 \u2212 \u03b5_\u03b8(x_t, t)\u2016\u00b2]  (predict noise)",
        "**Reverse:** x_{t-1} = (1/\u221a\u03b1_t)(x_t \u2212 \u03b2_t/\u221a(1-\u03b1\u0305_t) \u03b5_\u03b8) + \u03c3_t z",
        "",
        "Linear noise schedule: \u03b2_start=0.0001, \u03b2_end=0.1, T=100",
        "\u03b5-prediction (no tanh on output \u2014 predicts unbounded noise)",
    ], font_size=15)

    # Architecture
    s = make_content_slide(prs, "U-Net Architecture (~5.5M params)")
    try_add_image(s, 'fig8_architecture.png', Inches(0.3), Inches(1.0), width=Inches(9.4))

    # Training Config
    s = make_content_slide(prs, "Training Configuration")
    add_table(s, Inches(1.5), Inches(1.2), Inches(7), [
        ["Parameter", "Combined Model", "Original Model"],
        ["Demo data", "400 (combined)", "200 (left/right)"],
        ["Hidden dim", "256", "256"],
        ["Encoder blocks", "3", "6"],
        ["Batch size", "64", "256"],
        ["Learning rate", "2e-4", "1e-4"],
        ["EMA decay", "0.999", "0.999"],
        ["Epochs", "100", "500"],
    ], font_size=12)

    # DDIM
    s = make_content_slide(prs, "DDIM Inference (10-Step)")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5), [
        "100-step training \u2192 10-step DDIM inference (10\u00d7 speedup)",
        "\u03b7 = 0.3 (partial stochasticity for diversity)",
        "Temporal ensemble: weighted blend of overlapping predictions",
        "",
        "**Execution:** Predict 32-step chunk, execute first 8 steps, re-predict",
        "Provides both planning horizon (32) and responsiveness (replan every 8 steps)",
    ], font_size=16)

    # Training Results
    s = make_content_slide(prs, "Training Results")
    try_add_image(s, 'fig1_training_loss.png', Inches(0.3), Inches(1.0), width=Inches(5.5))
    add_bullet_list(s, Inches(6.0), Inches(1.5), Inches(3.5), Inches(2.5), [
        "**6,920 chunks** from 400 demos",
        "horizon=32, stride=2",
        "Loss: 0.154 \u2192 0.045",
    ], font_size=14)

    # Base Policy Eval
    s = make_content_slide(prs, "Base Policy Evaluation (50 Episodes)")
    add_table(s, Inches(0.5), Inches(1.2), Inches(4.5), [
        ["Metric", "Value"],
        ["Success Rate", "84% (42/50)"],
        ["Mean Episode Length", "344 steps"],
        ["Picked Left", "31/42 (73.8%)"],
        ["Picked Right", "11/42 (26.2%)"],
    ], font_size=13)
    add_bullet_list(s, Inches(5.5), Inches(1.2), Inches(4), Inches(3), [
        "Left-side bias reflects demo distribution",
        "16% failure rate \u2014 some episodes fail to lift block",
        "Policy learns multi-style distribution but doesn\u2019t distinguish at inference",
    ], font_size=14)

    # Debugging
    s = make_content_slide(prs, "Debugging Journey")
    add_table(s, Inches(0.5), Inches(1.2), Inches(9), [
        ["Problem", "Outcome", "Root Cause"],
        ["Initial eval", "13% success", "Checkpoint saved normalized stats (mean=0, std=1)"],
        ["Seed eval", "87% failure", "Modal collapse in some seeds"],
        ["Early architecture", "25% success", "Needed architecture refinements"],
    ], font_size=12)
    add_bullet_list(s, Inches(0.5), Inches(3.5), Inches(9), Inches(1.5), [
        "**Key lesson:** Normalization statistics must match inference conditions exactly"
    ], font_size=14)

    # Takeaways
    s = make_content_slide(prs, "Key Takeaways")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5), [
        "DDPM successfully learns multi-style trajectory distribution from mixed demos",
        "**84% task success** demonstrates viable base policy",
        "**But base policy is largely illegible** (VLO assessment in Week 3)",
        "Training-free steering is the key advantage \u2014 modify behavior without retraining",
        "Architecture is simple \u2014 1D U-Net with ResBlocks, not transformer-based",
    ], font_size=16)

    path = os.path.join(OUT_DIR, 'week2_diffusion_policy.pptx')
    prs.save(path)
    print(f"  \u2713 {os.path.basename(path)}")


# =================================================================
# WEEK 3: VLM Evaluation
# =================================================================
def build_week3():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    make_title_slide(prs,
        "Week 3: VLM-Based\nLegibility Evaluation",
        "Using Gemini 2.5 Flash as a Legibility Judge")

    # Evaluation Problem
    s = make_content_slide(prs, "The Evaluation Problem")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5), [
        "Human studies are gold standard but expensive and slow",
        "Analytical metrics (L_early) assume a specific Bayesian observer model",
        "**Proposal:** Use a Vision-Language Model (VLM) as a proxy human observer",
        "VLMs can watch trajectory videos and predict intent \u2014 just like a human",
        "",
        "**VLM Onset (VLO):** First timestep where VLM correctly predicts the robot\u2019s goal.",
        "Lower VLO = more legible.",
    ], font_size=16)

    # Pipeline
    s = make_content_slide(prs, "VLM Evaluation Pipeline")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5), [
        "1. Render policy rollout \u2192 MP4 video",
        "2. Extract **cumulative prefix frames**: k=6 windows over first 30% of episode",
        "3. For each window, send frames to Gemini 2.5 Flash with goal prompt",
        "4. VLM returns: {pA, pB, cue, choice} as structured JSON",
        "5. **VLO** = first k where choice matches true goal",
        "",
        "Cumulative windows (not independent) \u2014 VLM sees progressively more context",
    ], font_size=15)

    # Prompt Engineering
    s = make_content_slide(prs, "Prompt Engineering: V1 \u2192 V2")
    add_table(s, Inches(0.5), Inches(1.2), Inches(9), [
        ["Version", "Accuracy", "Key Issue"],
        ["V1", "45%", "World-coordinate confusion, long prompt, high temp"],
        ["V2", "97.5%", "Annotated reference frames, image-space reasoning, temp=0.1"],
    ], font_size=13)
    add_bullet_list(s, Inches(0.5), Inches(2.8), Inches(9), Inches(2.5), [
        "Annotated reference frame with block position markers (visual grounding)",
        "Image-space reasoning (\u201cLEFT/RIGHT in image\u201d) instead of world coordinates",
        "Short, direct prompt \u2014 removed verbose instructions",
        "Temperature lowered from default to 0.1",
    ], font_size=14)

    # Expert Demo Eval
    s = make_content_slide(prs, "Expert Demo Evaluation (Ground Truth)")
    add_table(s, Inches(0.5), Inches(1.2), Inches(4.5), [
        ["Metric", "Value"],
        ["Non-C accuracy", "94.7% (36/38)"],
        ["C-rate (uncertain)", "84.2% of windows"],
        ["Left arc accuracy", "95.0%"],
        ["Right arc accuracy", "100.0%"],
    ], font_size=12)
    add_table(s, Inches(5.5), Inches(1.2), Inches(4), [
        ["Style", "Mean VLO"],
        ["Legible", "2.93"],
        ["Neutral", "3.00"],
        ["Deceptive", "3.71"],
    ], font_size=13)
    add_text_box(s, Inches(0.5), Inches(4.2), Inches(9), Inches(0.5),
                 "VLO ordering: legible (2.93) < neutral (3.00) < deceptive (3.71) \u2713",
                 font_size=14, bold=True, color=GREEN)

    # VLO Distribution
    s = make_content_slide(prs, "VLO Distribution by Demo Style")
    try_add_image(s, 'fig6_vlo_distribution.png', Inches(0.3), Inches(1.0), width=Inches(9.0))

    # Base Policy VLO
    s = make_content_slide(prs, "Base Policy VLM Assessment (42 Episodes)")
    try_add_image(s, 'fig7_baseline_vlo.png', Inches(0.3), Inches(1.0), width=Inches(4.8))
    add_table(s, Inches(5.5), Inches(1.2), Inches(4), [
        ["Metric", "Value"],
        ["Mean VLO", "4.57 / 6"],
        ["Median VLO", "6.0 (worst)"],
        ["Never identified", "61.9%"],
        ["Legible at k=1", "9.5%"],
    ], font_size=12)

    # Key Insight
    s = make_content_slide(prs, "Key Insight")
    add_text_box(s, Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.0),
                 "The base policy achieves 84% task success\nbut produces largely illegible trajectories.",
                 font_size=22, bold=True, color=MAROON, alignment=PP_ALIGN.CENTER)
    add_bullet_list(s, Inches(0.8), Inches(3.2), Inches(8.4), Inches(2), [
        "Mixed training data averages out trajectory styles",
        "Diffusion policy doesn\u2019t prefer legible modes without explicit steering",
        "We need **training-free guidance** to push toward legible distribution region",
    ], font_size=16)

    # Limitations
    s = make_content_slide(prs, "What VLM Evaluation Does and Does NOT Show")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(4.2), Inches(3.5), [
        "**Supported:**",
        "\u2713 VLMs accurately predict goals (94.7%)",
        "\u2713 VLO correctly orders styles",
        "\u2713 Base policy is mostly illegible",
    ], font_size=14, color=GREEN)
    add_bullet_list(s, Inches(5.2), Inches(1.2), Inches(4.3), Inches(3.5), [
        "**Limitations:**",
        "\u2717 VLM \u2260 human observer",
        "\u2717 Post-hoc evaluation, not real-time",
        "\u2717 Only 10 block configs tested",
        "\u2717 VLO is a proxy metric",
    ], font_size=14, color=RED)

    path = os.path.join(OUT_DIR, 'week3_vlm_evaluation.pptx')
    prs.save(path)
    print(f"  \u2713 {os.path.basename(path)}")


# =================================================================
# WEEK 4: Guidance Methods
# =================================================================
def build_week4():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    make_title_slide(prs,
        "Week 4: Guidance Methods for\nLegible Diffusion",
        "Training-Free Steering of Trajectory Distribution")

    # The Guidance Idea
    s = make_content_slide(prs, "The Guidance Idea")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(2), [
        "Diffusion models generate via iterative denoising",
        "At each step, add gradient info to steer toward desired properties \u2014 **no retraining**",
    ], font_size=16)
    add_table(s, Inches(0.5), Inches(2.8), Inches(9), [
        ["#", "Method", "Type"],
        ["1", "Classifier Guidance", "Gradient-based, training-free"],
        ["2", "Best-of-N Reranking", "Sample + select, training-free"],
        ["3", "LegDiff (CFG)", "Conditional training required"],
        ["4", "VLM Text Reranking", "VLM as judge, training-free"],
    ], font_size=12)

    # Classifier Guidance
    s = make_content_slide(prs, "Method 1: Classifier Guidance")
    add_text_box(s, Inches(0.5), Inches(1.1), Inches(9), Inches(0.5),
                 "\u03b5\u0302 = \u03b5_\u03b8 \u2212 w\u221a(1\u2212\u03b1\u0305_t) \u2207 L_score(a_t)",
                 font_size=16, bold=True, color=MAROON)
    add_table(s, Inches(0.5), Inches(1.8), Inches(5), [
        ["Criterion", "Weight", "Description"],
        ["P_prox", "0.35", "Gaussian proximity to goal"],
        ["P_dir", "0.30", "Velocity alignment toward goal"],
        ["P_lat", "0.25", "Lateral separation from non-goal"],
        ["P_speed", "0.10", "Speed commitment"],
    ], font_size=11)
    add_text_box(s, Inches(0.5), Inches(4.2), Inches(9), Inches(0.4),
                 "Scoring function generated by Gemini from text prompt. Correlation with hand-crafted: r = 0.992",
                 font_size=13, color=GRAY)

    # Guidance Sweep
    s = make_content_slide(prs, "Classifier Guidance: Scale Sweep")
    try_add_image(s, 'fig2_guidance_sweep.png', Inches(0.3), Inches(1.0), width=Inches(5.0))
    add_table(s, Inches(5.5), Inches(1.2), Inches(4), [
        ["Scale w", "Success", "L_early"],
        ["0", "95%", "0.906"],
        ["5", "90%", "0.946"],
        ["10 *", "100%", "0.952"],
        ["20", "100%", "0.948"],
    ], font_size=12)
    add_text_box(s, Inches(5.5), Inches(4.2), Inches(4), Inches(0.4),
                 "* w=10 optimal", font_size=13, bold=True, color=MAROON)

    # Terminology
    s = make_content_slide(prs, "Terminology Correction: NOT DPS")
    add_table(s, Inches(0.5), Inches(1.2), Inches(9), [
        ["", "Classifier Guidance (Ours)", "True DPS (Chung 2023)"],
        ["Gradient applied to", "Noise prediction \u03b5_\u03b8", "Denoised sample x_{t-1}"],
        ["Our result", "100% success", "0% success"],
        ["Step divergence", "0%", "38%"],
    ], font_size=12)
    add_text_box(s, Inches(0.5), Inches(3.5), Inches(9), Inches(0.5),
                 "All reported results use classifier guidance (Dhariwal & Nichol, 2021), NOT DPS.",
                 font_size=14, bold=True, color=RED)

    # Best-of-N
    s = make_content_slide(prs, "Method 2: Best-of-N Reranking")
    try_add_image(s, 'fig3_best_of_n.png', Inches(0.3), Inches(1.0), width=Inches(5.0))
    add_table(s, Inches(5.5), Inches(1.2), Inches(4), [
        ["N", "L_early", "\u0394 vs N=1"],
        ["1", "0.732", "\u2014"],
        ["4", "0.779", "+6.4%"],
        ["8", "0.797", "+8.8%"],
        ["16", "0.806", "+10.1%"],
    ], font_size=12)
    add_text_box(s, Inches(5.5), Inches(4.2), Inches(4), Inches(0.4),
                 "Simple, gradient-free, diminishing returns above N=16",
                 font_size=12, color=GRAY)

    # LegDiff
    s = make_content_slide(prs, "Method 3: LegDiff (Classifier-Free Guidance)")
    add_text_box(s, Inches(0.5), Inches(1.1), Inches(9), Inches(0.5),
                 "\u03b5\u0302 = \u03b5_\u03b8(a_t, \u2205) + w \u00b7 (\u03b5_\u03b8(a_t, G) \u2212 \u03b5_\u03b8(a_t, \u2205))",
                 font_size=16, bold=True, color=MAROON)
    add_bullet_list(s, Inches(0.5), Inches(1.8), Inches(9), Inches(1.5), [
        "Trains both conditional (goal-aware) and unconditional models",
        "Conv1d backbone with temporal convolutions (kernel=5), w=3",
    ], font_size=14)
    add_table(s, Inches(2), Inches(3.2), Inches(6), [
        ["Condition", "Success", "L_early"],
        ["Baseline", "100%", "0.922"],
        ["LegDiff (w=3)", "100%", "0.935"],
    ], font_size=13)

    # VLM Reranking
    s = make_content_slide(prs, "Method 4: VLM Text-Based Reranking")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(1.5), [
        "1. Generate K=5 candidates using classifier guidance",
        "2. Render each as video \u2192 VLM scores for goal identifiability",
        "3. Execute the trajectory VLM most confidently identifies",
    ], font_size=14)
    add_table(s, Inches(1.5), Inches(3.0), Inches(7), [
        ["Condition", "Success", "L_early"],
        ["Single candidate", "85%", "0.946"],
        ["Oracle rerank", "95%", "0.968"],
        ["VLM rerank", "100%", "0.972"],
    ], font_size=13)

    # All Methods
    s = make_content_slide(prs, "All Methods Compared")
    try_add_image(s, 'fig4_method_comparison.png', Inches(0.3), Inches(1.0), width=Inches(9.0))

    # VLM-Generated scoring
    s = make_content_slide(prs, "Honest: VLM-Generated Scoring Function")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5), [
        "Generated by Gemini from **text-only** prompt (no visual data)",
        "Single-shot, no iterative refinement (not full EUREKA protocol)",
        "Achieves r=0.992 correlation with hand-crafted 4-criteria function",
        "**Any capable text LLM could produce equivalent output**",
        "",
        "Claimed as: valid use of LLM code generation, NOT a visual contribution",
    ], font_size=16)

    path = os.path.join(OUT_DIR, 'week4_guidance_methods.pptx')
    prs.save(path)
    print(f"  \u2713 {os.path.basename(path)}")


# =================================================================
# WEEK 5: Full Pipeline
# =================================================================
def build_week5():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    make_title_slide(prs,
        "Week 5: Full Pipeline Integration\n& Lessons Learned",
        "From Training to VLM-Guided Legible Diffusion")

    # Complete Pipeline
    s = make_content_slide(prs, "The Complete Pipeline")
    add_bullet_list(s, Inches(1), Inches(1.2), Inches(8), Inches(3.8), [
        "1. Demo Collection (400 B\u00e9zier trajectories, 3 styles)",
        "2. Diffusion Policy Training (DDPM, 100 epochs, 5.5M params)",
        "3. Base Policy: 84% success, VLO=4.57 \u2014 illegible",
        "4. Classifier Guidance (w=10): L_early=0.952, 100% success",
        "5. VLM Reranking (K=5): L_early=0.972, 100% success",
        "6. VLM Evaluation (VLO metric, Gemini 2.5 Flash)",
    ], font_size=16)

    # What Worked
    s = make_content_slide(prs, "What Worked")
    add_table(s, Inches(0.3), Inches(1.2), Inches(9.4), [
        ["Component", "Result", "Why It Worked"],
        ["Mixed-style demos", "84% success", "B\u00e9zier gives precise control; 400 demos sufficient"],
        ["DDPM U-Net", "Learns distribution", "Suitable for multimodal action data"],
        ["Classifier guidance", "+5.2% L_early", "Gradient signal meaningful; w=10 balances task/legibility"],
        ["VLM evaluation", "94.7% accuracy", "Prompt engineering (v1\u2192v2) was critical"],
        ["VLM reranking", "Best L_early", "VLM as judge catches what metrics miss"],
    ], font_size=11)

    # What Didn't Work
    s = make_content_slide(prs, "What Didn't Work")
    add_table(s, Inches(0.3), Inches(1.2), Inches(9.4), [
        ["Attempt", "Outcome", "Root Cause"],
        ["True DPS", "0% success", "Gradient to x_{t-1} causes divergence"],
        ["Initial eval", "13% success", "Checkpoint bug: wrong normalized stats"],
        ["VLM v1 prompt", "45% accuracy", "World-coordinate confusion"],
        ["Seed eval (10 seeds)", "87% failure", "Modal collapse in some seeds"],
    ], font_size=11)

    # Engineering Lessons
    s = make_content_slide(prs, "Engineering Lessons")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5), [
        "**Video saving bug** \u2014 reset() cleared _video_frames; fix: guard with flag",
        "**Video recording order** \u2014 record_video() before reset() erased initial frame",
        "**Camera convention** \u2014 Image-left \u2260 world-left caused systematic VLM errors",
        "**Normalization stats** \u2014 Wrong stats in checkpoint \u2192 wildly wrong predictions",
    ], font_size=15)

    # Pipeline Stages
    s = make_content_slide(prs, "Full Pipeline: Staged Improvement")
    try_add_image(s, 'fig5_pipeline_stages.png', Inches(0.3), Inches(1.0), width=Inches(5.5))
    add_table(s, Inches(6.0), Inches(1.2), Inches(3.5), [
        ["Stage", "L_early"],
        ["0: Baseline", "0.898"],
        ["1: +Guidance", "0.937"],
        ["2: +VLM", "0.972"],
    ], font_size=12)
    add_text_box(s, Inches(6.0), Inches(3.8), Inches(3.5), Inches(0.5),
                 "p = 0.00042", font_size=14, bold=True, color=MAROON)

    # Honest Assessment
    s = make_content_slide(prs, "Honest Assessment (6 Issues)")
    add_table(s, Inches(0.3), Inches(1.2), Inches(9.4), [
        ["#", "Issue", "Severity"],
        ["1", "Method is classifier guidance, NOT DPS", "High"],
        ["2", "VLM scoring from text only (not visual)", "Medium"],
        ["3", "Single-shot, not full EUREKA", "Medium"],
        ["4", "No human study to validate VLM", "High"],
        ["5", "No reverse-steering (w<0) test", "Medium"],
        ["6", "Same configs for train and eval", "Medium"],
    ], font_size=11)

    # What We CAN / CANNOT Claim
    s = make_content_slide(prs, "Supported vs. Unsupported Claims")
    add_bullet_list(s, Inches(0.3), Inches(1.2), Inches(4.4), Inches(3.5), [
        "**CAN Claim (with evidence):**",
        "\u2713 Multi-style learning: 84% success",
        "\u2713 VLM accuracy: 94.7%",
        "\u2713 VLO orders styles correctly",
        "\u2713 Guidance significant: p=0.00042",
        "\u2713 Pipeline: 0.898 \u2192 0.972",
    ], font_size=13, color=GREEN)
    add_bullet_list(s, Inches(5.0), Inches(1.2), Inches(4.5), Inches(3.5), [
        "**CANNOT Claim:**",
        "\u2717 More legible to humans",
        "\u2717 VLM scoring > hand-crafted",
        "\u2717 Generalizes to real robots",
        "\u2717 DPS inherently unsuitable",
        "\u2717 Visual contribution from VLM",
    ], font_size=13, color=RED)

    # Future
    s = make_content_slide(prs, "Future Directions")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5), [
        "1. **Human validation study** \u2014 Does VLO correspond to faster human prediction?",
        "2. **Online VLM guidance** \u2014 Real-time VLM feedback during execution",
        "3. **Reverse-steering** \u2014 Test w < 0 for causal evidence",
        "4. **Real robot transfer** \u2014 Franka Panda physical experiments",
        "5. **Multi-goal settings** \u2014 3+ objects for richer legibility challenges",
        "6. **Full EUREKA loop** \u2014 Iterative refinement with environment feedback",
    ], font_size=15)

    # Thank You
    s = make_section_slide(prs, "Thank You", "Questions?")

    path = os.path.join(OUT_DIR, 'week5_full_pipeline.pptx')
    prs.save(path)
    print(f"  \u2713 {os.path.basename(path)}")


# =================================================================
# FINAL DEFENSE
# =================================================================
def build_defense():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Slide 1: Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, MAROON)
    add_shape_bg(s, Inches(0), Inches(3.4), Inches(10), Inches(0.08), GOLD)
    add_text_box(s, Inches(0.5), Inches(0.8), Inches(9), Inches(1.8),
                 "VLM-Guided Diffusion Policies\nfor Legible Robot Motion\nin Multi-Goal Manipulation",
                 font_size=30, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(2.8), Inches(9), Inches(0.5),
                 "Anudeep Gottapu", font_size=18, color=GOLD, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(3.7), Inches(9), Inches(0.5),
                 "MS Thesis Defense  |  Arizona State University  |  Spring 2026",
                 font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(4.3), Inches(9), Inches(0.5),
                 "Committee: [Advisor Name], [Member 2], [Member 3]",
                 font_size=12, color=RGBColor(0xCC, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)

    # Slide 2: Motivation
    s = make_content_slide(prs, "Motivation: Legibility in Human-Robot Collaboration")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5), [
        "In shared workspaces, robots must communicate intent through motion",
        "**Legible motion** = observer quickly infers the goal (Dragan et al., 2013)",
        "Ambiguous trajectories \u2192 human waits \u2192 slower, less safe collaboration",
        "",
        "**This work:** Can diffusion policies generate **both successful AND legible**",
        "manipulation trajectories, guided by Vision-Language Model feedback?",
    ], font_size=16)

    # Slide 3: Contributions
    s = make_content_slide(prs, "Contributions")
    add_bullet_list(s, Inches(0.5), Inches(1.3), Inches(9), Inches(3.5), [
        "**1. VLM-based legibility evaluation pipeline**",
        "   Gemini 2.5 Flash achieves 94.7% accuracy. Defines VLM Onset (VLO) metric.",
        "",
        "**2. Systematic comparison of 4 training-free guidance methods**",
        "   Classifier guidance, Best-of-N, Classifier-Free Guidance, VLM reranking.",
        "",
        "**3. End-to-end pipeline: 0.898 \u2192 0.972 L_early (+7.4pp)**",
        "   At 100% task success, p = 0.00042 (statistically significant).",
    ], font_size=15)

    # Slide 4: Environment
    s = make_content_slide(prs, "Environment: TwoBlockPick")
    add_table(s, Inches(0.5), Inches(1.2), Inches(5), [
        ["Parameter", "Value"],
        ["Robot", "Franka Panda 7-DOF (PyBullet)"],
        ["Observation", "22-dim (ee + grip + blocks)"],
        ["Action", "5-dim (\u0394x, \u0394y, \u0394z, \u0394yaw, grip)"],
        ["Goal", "Pick either block, lift above z=0.52m"],
        ["Blocks", "Two red, 4cm, symmetric"],
    ], font_size=12)
    add_text_box(s, Inches(6), Inches(1.2), Inches(3.5), Inches(1),
                 "Camera: yaw=135\u00b0\nImage-left = World-right",
                 font_size=14, bold=True, color=MAROON)

    # Slide 5: Demonstrations
    s = make_content_slide(prs, "400 Expert Demos via B\u00e9zier Trajectories")
    add_table(s, Inches(0.5), Inches(1.2), Inches(9), [
        ["Style", "Count (%)", "Design"],
        ["Legible", "200 (50%)", "Control point toward goal. Commitment by t\u22480.3"],
        ["Neutral", "100 (25%)", "Control point at y=0. No early signal"],
        ["Deceptive", "100 (25%)", "Feint wrong side, then commit to correct goal"],
    ], font_size=12)
    add_text_box(s, Inches(0.5), Inches(3.3), Inches(9), Inches(0.5),
                 "10 configurations \u00d7 40 demos/config | Purpose: train on diverse styles, evaluate guidance",
                 font_size=14, color=GRAY)

    # Slide 6: Architecture
    s = make_content_slide(prs, "Diffusion Policy: DDPM U-Net (~5.5M params)")
    try_add_image(s, 'fig8_architecture.png', Inches(0.2), Inches(1.0), width=Inches(9.5))

    # Slide 7: Base Policy
    s = make_content_slide(prs, "Base Policy Results")
    add_table(s, Inches(0.5), Inches(1.2), Inches(4), [
        ["Metric", "Value"],
        ["Success Rate", "84% (42/50)"],
        ["Mean Steps", "344"],
    ], font_size=13)
    add_table(s, Inches(5), Inches(1.2), Inches(4.5), [
        ["VLM Assessment", "Value"],
        ["Mean VLO", "4.57 / 6"],
        ["Median VLO", "6.0 (worst)"],
        ["Never identified", "61.9%"],
    ], font_size=13)
    add_text_box(s, Inches(0.5), Inches(3.8), Inches(9), Inches(0.5),
                 "Finding: 84% success but mostly illegible. Need guidance.",
                 font_size=16, bold=True, color=MAROON, alignment=PP_ALIGN.CENTER)

    # Slide 8: VLM Pipeline
    s = make_content_slide(prs, "VLM Evaluation Pipeline")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(5.5), Inches(3.5), [
        "1. Render episode \u2192 MP4 video",
        "2. Extract k=6 cumulative prefix windows (first 30%)",
        "3. VLM predicts goal: {pA, pB, cue, choice}",
        "4. **VLO** = first window with correct prediction",
        "",
        "Prompt V1: 45% accuracy",
        "Prompt **V2: 97.5% accuracy**",
    ], font_size=14)
    add_table(s, Inches(6.5), Inches(1.2), Inches(3), [
        ["Style", "Mean VLO"],
        ["Legible", "2.93"],
        ["Neutral", "3.00"],
        ["Deceptive", "3.71"],
    ], font_size=12)

    # Slide 9: Classifier Guidance
    s = make_content_slide(prs, "Method 1: Classifier Guidance")
    try_add_image(s, 'fig2_guidance_sweep.png', Inches(0.2), Inches(1.0), width=Inches(5.2))
    add_bullet_list(s, Inches(5.5), Inches(1.2), Inches(4), Inches(3), [
        "\u03b5\u0302 = \u03b5_\u03b8 \u2212 w\u221a(1\u2212\u03b1\u0305_t) \u2207L",
        "**w=10 optimal**",
        "100% success, L_early=0.952",
        "",
        "Note: This is classifier guidance",
        "(Dhariwal & Nichol, 2021)",
        "NOT DPS. True DPS \u2192 0% success",
    ], font_size=13)

    # Slide 10: Best-of-N
    s = make_content_slide(prs, "Method 2: Best-of-N Reranking")
    try_add_image(s, 'fig3_best_of_n.png', Inches(0.2), Inches(1.0), width=Inches(5.2))
    add_table(s, Inches(5.8), Inches(1.2), Inches(3.8), [
        ["N", "L_early", "\u0394"],
        ["1", "0.732", "\u2014"],
        ["4", "0.779", "+6.4%"],
        ["8", "0.797", "+8.8%"],
        ["16", "0.806", "+10.1%"],
    ], font_size=12)

    # Slide 11: LegDiff
    s = make_content_slide(prs, "Method 3: LegDiff (Classifier-Free Guidance)")
    add_text_box(s, Inches(0.5), Inches(1.1), Inches(9), Inches(0.5),
                 "\u03b5\u0302 = \u03b5_\u03b8(a_t, \u2205) + w \u00b7 (\u03b5_\u03b8(a_t, G) \u2212 \u03b5_\u03b8(a_t, \u2205))",
                 font_size=16, bold=True, color=MAROON)
    add_table(s, Inches(2), Inches(2.0), Inches(6), [
        ["Condition", "Success", "L_early"],
        ["Baseline", "100%", "0.922"],
        ["LegDiff (w=3)", "100%", "0.935"],
    ], font_size=14)
    add_text_box(s, Inches(0.5), Inches(3.8), Inches(9), Inches(0.5),
                 "Conv1d backbone, conditional + unconditional training. Modest but consistent improvement.",
                 font_size=13, color=GRAY)

    # Slide 12: VLM Reranking
    s = make_content_slide(prs, "Method 4: VLM Text Reranking")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(5), Inches(2), [
        "Generate K=5 candidates with classifier guidance",
        "Render each as video",
        "VLM selects most identifiable trajectory",
    ], font_size=14)
    add_table(s, Inches(5.5), Inches(1.2), Inches(4), [
        ["Method", "Success", "L_early"],
        ["Single", "85%", "0.946"],
        ["Oracle", "95%", "0.968"],
        ["VLM rerank", "100%", "0.972"],
    ], font_size=13)

    # Slide 13: Pipeline Stages
    s = make_content_slide(prs, "Full Pipeline: Staged Results")
    try_add_image(s, 'fig5_pipeline_stages.png', Inches(0.2), Inches(1.0), width=Inches(5.5))
    add_table(s, Inches(6.0), Inches(1.2), Inches(3.5), [
        ["Stage", "Success", "L_early"],
        ["0: Baseline", "80%", "0.898"],
        ["1: +Guidance", "100%", "0.937"],
        ["2: +VLM", "100%", "0.972"],
    ], font_size=12)
    add_text_box(s, Inches(6.0), Inches(4.0), Inches(3.5), Inches(0.5),
                 "\u0394 = +0.074, p = 0.00042", font_size=14, bold=True, color=MAROON)

    # Slide 14: All Methods
    s = make_content_slide(prs, "All Methods Compared")
    try_add_image(s, 'fig4_method_comparison.png', Inches(0.2), Inches(1.0), width=Inches(9.5))

    # Slide 15: Honest Assessment
    s = make_content_slide(prs, "Honest Assessment")
    add_table(s, Inches(0.3), Inches(1.2), Inches(9.4), [
        ["Issue", "Severity", "Detail"],
        ["Terminology", "High", "Method is classifier guidance, NOT DPS"],
        ["Text-only scoring", "Medium", "VLM never saw visual data; any LLM could produce similar"],
        ["Not full EUREKA", "Medium", "Single-shot, no iterative refinement"],
        ["No human study", "High", "VLM \u2260 human observer; correlation assumed"],
        ["No reverse-steering", "Medium", "w<0 never tested; no causal proof"],
        ["Same train/test configs", "Medium", "Generalization not demonstrated"],
    ], font_size=11)

    # Slide 16: Supported Claims
    s = make_content_slide(prs, "Supported Claims (All Evidence-Backed)")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5), [
        "\u2713 Diffusion policies learn multi-style trajectories (84% success \u2014 Table 5.1)",
        "\u2713 VLMs accurately evaluate legibility (94.7% accuracy \u2014 Table 5.3)",
        "\u2713 VLO correctly orders styles: legible < neutral < deceptive (Table 5.3)",
        "\u2713 Classifier guidance improves legibility significantly (p=0.00042 \u2014 Table 5.7)",
        "\u2713 Full pipeline: 0.898 \u2192 0.972 at 100% success (Table 5.7)",
    ], font_size=15, color=GREEN)

    # Slide 17: Future Work
    s = make_content_slide(prs, "Future Work")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5), [
        "1. **Human study** \u2014 Validate VLO correlates with human prediction speed",
        "2. **Online VLM guidance** \u2014 Real-time VLM feedback during execution",
        "3. **Reverse-steering** \u2014 Test w<0 for causal evidence",
        "4. **Real robot** \u2014 Franka Panda physical experiments",
        "5. **Multi-goal** \u2014 3+ objects for richer legibility challenges",
        "6. **Full EUREKA** \u2014 Iterative refinement with environment feedback",
    ], font_size=15)

    # Slide 18: Summary
    s = make_content_slide(prs, "Summary")
    add_text_box(s, Inches(0.5), Inches(1.1), Inches(9), Inches(0.8),
                 "Training-free guidance over diffusion policies is a viable and effective\napproach to legible motion generation in multi-goal manipulation.",
                 font_size=16, bold=True, color=MAROON, alignment=PP_ALIGN.CENTER)
    add_bullet_list(s, Inches(0.5), Inches(2.2), Inches(9), Inches(2.5), [
        "VLM evaluation pipeline: 94.7% accuracy, VLO metric validated",
        "Classifier guidance at w=10: 100% success, L_early = 0.952",
        "Full pipeline with VLM reranking: L_early = 0.972 (p=0.00042 vs baseline)",
        "4 guidance methods compared on same task, same evaluation framework",
        "",
        "**Honest limitations:** No human study, same configs, terminology corrections needed.",
    ], font_size=14)

    # Slide 19: Thank You
    s = make_section_slide(prs, "Thank You", "Questions?")
    add_text_box(s, Inches(0.8), Inches(3.5), Inches(8.4), Inches(0.5),
                 "Anudeep Gottapu  |  Arizona State University",
                 font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

    path = os.path.join(OUT_DIR, 'final_defense.pptx')
    prs.save(path)
    print(f"  \u2713 {os.path.basename(path)}")


# =================================================================
if __name__ == '__main__':
    print('Generating PowerPoint presentations...')
    build_week1()
    build_week2()
    build_week3()
    build_week4()
    build_week5()
    build_defense()
    print(f'\nAll .pptx files saved to: {OUT_DIR}')
