#!/usr/bin/env python
"""
Build a comprehensive progress presentation for professor review.

Covers every decision point and architectural change from Feb 8 – Mar 31 2026.
Embeds figures, video composites, tables with real numbers from result JSONs.

Output:  presentation/thesis_progress_v2.pptx
"""

import io
import json
import sys
from pathlib import Path

import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent.parent
FIGURES = ROOT / "figures"
DEMO_VIDS = ROOT / "data" / "demos" / "demo_videos"
OUTPUTS = ROOT / "outputs"
RUNS = ROOT / "runs"
OUT_DIR = ROOT / "presentation"
OUT_DIR.mkdir(exist_ok=True)
PPTX_PATH = OUT_DIR / "thesis_progress_v2.pptx"

# ── Colours ────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1B, 0x2A, 0x4A)
MED_BLUE    = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE  = RGBColor(0xDB, 0xEA, 0xFE)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
ORANGE      = RGBColor(0xEA, 0x58, 0x0C)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
RED         = RGBColor(0xDC, 0x26, 0x26)
GRAY        = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY  = RGBColor(0xF3, 0xF4, 0xF6)
GOLD        = RGBColor(0xF5, 0x9E, 0x0B)
LIGHT_GREEN = RGBColor(0xEC, 0xFD, 0xF5)
LIGHT_RED   = RGBColor(0xFE, 0xE2, 0xE2)
LIGHT_GOLD  = RGBColor(0xFF, 0xFB, 0xEB)
DARK_GREEN  = RGBColor(0x06, 0x6B, 0x52)

# ═════════════════════════════════════════════════════════════════════════
# UTILITY FUNCTIONS (reused from build_presentation.py)
# ═════════════════════════════════════════════════════════════════════════

def extract_frame(video_path: Path, frame_idx: int = 30) -> bytes:
    cap = cv2.VideoCapture(str(video_path))
    cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
    ret, frame = cap.read()
    cap.release()
    if not ret:
        img = Image.new("RGB", (640, 480), (200, 200, 200))
        buf = io.BytesIO(); img.save(buf, format="PNG"); return buf.getvalue()
    frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
    img = Image.fromarray(frame_rgb)
    buf = io.BytesIO(); img.save(buf, format="PNG"); return buf.getvalue()


def extract_composite_frames(video_path: Path, n_frames: int = 6) -> bytes:
    cap = cv2.VideoCapture(str(video_path))
    total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    if total <= 0:
        cap.release()
        img = Image.new("RGB", (640, 480), (200, 200, 200))
        buf = io.BytesIO(); img.save(buf, format="PNG"); return buf.getvalue()
    indices = np.linspace(0, total - 1, n_frames, dtype=int)
    frames = []
    for idx in indices:
        cap.set(cv2.CAP_PROP_POS_FRAMES, int(idx))
        ret, frame = cap.read()
        if ret:
            frames.append(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
    cap.release()
    if not frames:
        img = Image.new("RGB", (640, 480), (200, 200, 200))
        buf = io.BytesIO(); img.save(buf, format="PNG"); return buf.getvalue()
    target_h = 240
    resized = []
    for f in frames:
        h, w = f.shape[:2]
        scale = target_h / h
        new_w = int(w * scale)
        resized.append(cv2.resize(f, (new_w, target_h)))
    composite = np.concatenate(resized, axis=1)
    img = Image.fromarray(composite)
    buf = io.BytesIO(); img.save(buf, format="PNG"); return buf.getvalue()


def set_slide_bg(slide, color):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(font_size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font_name; p.alignment = alignment
    return txBox


def add_bullet_slide(slide, left, top, width, height, bullets,
                     font_size=16, color=BLACK, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet; p.font.size = Pt(font_size)
        p.font.color.rgb = color; p.font.name = "Calibri"
        p.space_after = spacing; p.level = 0
    return txBox


def add_title_bar(slide, title_text, subtitle_text=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    shape.fill.solid(); shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
                title_text, font_size=28, bold=True, color=WHITE)
    if subtitle_text:
        add_textbox(slide, Inches(0.5), Inches(0.65), Inches(12), Inches(0.35),
                    subtitle_text, font_size=14, color=LIGHT_BLUE)


def add_table(slide, left, top, width, height, data, col_widths=None):
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r_idx, row in enumerate(data):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE
            elif r_idx % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY
    return table_shape


def add_phase_badge(slide, phase_num, date_range, x=11.0, y=0.2):
    """Small badge in top-right showing phase number and date."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(2.0), Inches(0.7))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Phase {phase_num}\n{date_range}"
    p.font.size = Pt(9); p.font.bold = True
    p.font.color.rgb = DARK_BLUE; p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER


def try_add_figure(slide, fig_name, left, top, width, height):
    """Try to add a figure; return True if found."""
    fig_path = FIGURES / fig_name
    if fig_path.exists():
        slide.shapes.add_picture(str(fig_path), left, top, width, height)
        return True
    return False


def try_add_video_frame(slide, vid_name, left, top, width, height, frame_idx=30):
    """Try to add a frame extracted from a demo video."""
    vid_path = DEMO_VIDS / vid_name
    if vid_path.exists():
        frame_bytes = extract_frame(vid_path, frame_idx)
        slide.shapes.add_picture(io.BytesIO(frame_bytes), left, top, width, height)
        return True
    return False


def try_add_video_composite(slide, vid_name, left, top, width, height, n_frames=6):
    """Try to add a composite strip from a demo video."""
    vid_path = DEMO_VIDS / vid_name
    if vid_path.exists():
        comp = extract_composite_frames(vid_path, n_frames)
        slide.shapes.add_picture(io.BytesIO(comp), left, top, width, height)
        return True
    return False


# ═════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS  (25 slides, chronological)
# ═════════════════════════════════════════════════════════════════════════

# ── Slide 1: Title ──────────────────────────────────────────────────────
def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BLUE)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
                "Legible Robot Trajectories via\nDiffusion Policy + VLM Guidance",
                font_size=38, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(0.8),
                "Thesis Progress Report: Decisions, Debugging & Results\n"
                "TwoBlockPick  |  Franka Panda  |  Gemini VLM  |  LegDiff CFG",
                font_size=16, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.1), Inches(5.333), Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

    add_textbox(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.5),
                "February 8 \u2013 March 31, 2026  |  10 Engineering Phases",
                font_size=14, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(5.3), Inches(11.3), Inches(0.5),
                "April 2026",
                font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)


# ── Slide 2: Research Context ──────────────────────────────────────────
def slide_02_research_context(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Research Context & Goal",
                  "What is legibility and why does it matter for human-robot interaction?")

    bullets_left = [
        "Legibility (Dragan et al., 2013): An observer can predict the\n"
        "robot's goal early in the trajectory (first 30-40% of motion).",
        "",
        "Problem: Diffusion policies produce efficient but ambiguous\n"
        "straight-line paths \u2014 an observer cannot tell which block the\n"
        "robot intends to pick until it is almost there.",
        "",
        "Goal: Steer the policy at inference time to produce curved\n"
        "(arcing) trajectories that reveal intent earlier, without\n"
        "sacrificing task success and without retraining.",
        "",
        "Four steering setups explored:",
        "  1. Online VLM Best-of-K reranking (Feb 24)",
        "  2. Proxy geometric scorer (Feb 26)",
        "  3. Paired replanning with VLM (Feb 28)",
        "  4. LegDiff: CFG-conditioned diffusion (Mar 30)",
    ]
    add_bullet_slide(slide, Inches(0.5), Inches(1.3), Inches(6.8), Inches(6.0),
                     bullets_left, font_size=13, color=BLACK)

    # Environment image
    try_add_video_frame(slide, "cfg00_left_arc00.mp4",
                        Inches(7.8), Inches(1.4), Inches(5.0), Inches(3.7), 0)
    add_textbox(slide, Inches(7.8), Inches(5.2), Inches(5.0), Inches(0.8),
                "TwoBlockPick: Franka Panda, two red blocks\n"
                "at y=\u00b10.07m, 22-dim obs, 5-dim actions",
                font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)


# ── Slide 3: Task & Data ──────────────────────────────────────────────
def slide_03_task_and_data(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Task, Environment & Demonstration Data",
                  "400 B\u00e9zier demonstrations  |  20 arc levels \u00d7 2 sides \u00d7 10 configs")
    add_phase_badge(slide, 1, "Feb 8\u201313")

    # Three environment frames at different arc levels
    videos = [
        ("cfg00_left_arc00.mp4", "Arc 00 \u2014 Straight"),
        ("cfg00_left_arc10.mp4", "Arc 10 \u2014 Moderate"),
        ("cfg00_left_arc19.mp4", "Arc 19 \u2014 Strong"),
    ]
    x_positions = [0.3, 4.5, 8.7]
    for (vname, label), xp in zip(videos, x_positions):
        try_add_video_frame(slide, vname, Inches(xp), Inches(1.3),
                            Inches(4.0), Inches(2.8), 60)
        add_textbox(slide, Inches(xp), Inches(4.15), Inches(4.0), Inches(0.3),
                    label, font_size=12, bold=True, color=DARK_BLUE,
                    alignment=PP_ALIGN.CENTER)

    specs = [
        "Observation 22-dim: EE pose (7) + block pos (6) + gripper (2) + joints (7)",
        "Action 5-dim: \u0394x, \u0394y, \u0394z, \u0394yaw, gripper  |  Scale: 0.05 m/step",
        "Camera: yaw=135\u00b0, pitch=\u221230\u00b0, dist=0.9m  |  640\u00d7480 @ 30Hz",
        "B\u00e9zier arc index 0\u201319 controls curvature; class 15\u201319 = legible arcs (arc > 0.134m)",
        "Each config places blocks at fixed positions; arc determines approach trajectory shape",
    ]
    add_bullet_slide(slide, Inches(0.4), Inches(4.7), Inches(12.5), Inches(2.5),
                     specs, font_size=12, color=GRAY)


# ── Slide 4: Architecture Search ─────────────────────────────────────
def slide_04_architecture_search(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Decision 1: Architecture Search \u2014 12 Training Runs",
                  "Why we settled on UNet Conv1d (8.7M params) over smaller MLP models")
    add_phase_badge(slide, 2, "Feb 8\u201322")

    data = [
        ["Run", "Date", "Architecture", "Params", "Epochs", "Success", "Decision"],
        ["1", "Feb 8", "MLP 3-layer", "~1M", "300", "0%", "Too small"],
        ["2", "Feb 10", "MLP 6-layer", "~3M", "300", "0%", "Still fails"],
        ["3", "Feb 11", "Conv1d shallow", "~5M", "300", "0%", "Under capacity"],
        ["4-7", "Feb 11\u201313", "Various MLPs", "1\u20135M", "300", "0%", "All fail"],
        ["8", "Feb 13", "UNet Conv1d", "8.7M", "300", "14%", "\u2714 First success"],
        ["9", "Feb 17", "UNet H=48 n=16", "8.7M", "300", "0\u20138%", "Horizon broke it"],
        ["10", "Feb 22", "UNet + bug fixes", "8.7M", "100", "84\u201396%", "\u2714 FINAL"],
    ]
    add_table(slide, Inches(0.3), Inches(1.3), Inches(12.7), Inches(3.2), data)

    bullets = [
        "Key insight: Small models (1\u20135M params) cannot represent the 400-demo multimodal "
        "distribution with 20 arc levels. UNet Conv1d with skip connections was necessary.",
        "Run 8 was the first to succeed \u2014 14% at epoch 300. Runs 1\u20137 all got 0%.",
        "Run 9 changed horizon 32\u219248 and execute_steps 8\u219216 simultaneously \u2014 broke the "
        "policy entirely. Reverted in Run 10.",
        "Run 10 with critical bug fixes (see next slides) achieved 84\u201396% success.",
    ]
    add_bullet_slide(slide, Inches(0.4), Inches(4.8), Inches(12.5), Inches(2.5),
                     bullets, font_size=12, color=BLACK)


# ── Slide 5: Bug #1 \u2014 execute_steps ──────────────────────────────────
def slide_05_execute_steps_bug(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Decision 2: Fixing execute_steps \u2014 0% \u2192 50% Success",
                  "Training loss was 0.004 (good) but simulation success was 0% (bad)")
    add_phase_badge(slide, 3, "Feb 16")

    bullets_left = [
        "Symptom: Checkpoint loss was excellent (0.004) but the\n"
        "robot failed every episode in simulation.",
        "",
        "Root Cause: execute_steps was set to 16 in eval,\n"
        "but the model was trained on consecutive timesteps.",
        "  \u2022 Training sees: obs[t], obs[t+1], obs[t+2], ...",
        "  \u2022 With es=16: obs[t], obs[t+16], obs[t+32], ...",
        "  \u2022 The model sees OOD (out-of-distribution) observations",
        "  \u2022 \u2192 erratic actions \u2192 0% success",
        "",
        "Fix: Set execute_steps = 8 everywhere.",
        "  \u2022 es=1: 0% (too chaotic, constant replanning)",
        "  \u2022 es=8: 10\u201350% \u2714 (matches training temporal structure)",
        "  \u2022 es=16: 0% \u2718 (OOD observations)",
        "",
        "Lesson: Eval temporal structure must match training.",
    ]
    add_bullet_slide(slide, Inches(0.4), Inches(1.3), Inches(7.0), Inches(6.0),
                     bullets_left, font_size=13, color=BLACK)

    # Table on the right
    es_data = [
        ["execute_steps", "Success", "Status"],
        ["1", "0%", "\u2718 Chaotic"],
        ["8", "10\u201350%", "\u2714 Optimal"],
        ["16", "0%", "\u2718 OOD obs"],
    ]
    add_table(slide, Inches(8.0), Inches(1.5), Inches(4.8), Inches(2.5), es_data)

    add_textbox(slide, Inches(8.0), Inches(4.5), Inches(4.8), Inches(1.5),
                "This was the first major debugging\n"
                "breakthrough. It took 3 days of\n"
                "systematic testing to isolate\n"
                "execute_steps as the root cause.",
                font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)


# ── Slide 6: Bug #2 \u2014 Horizon Mismatch ────────────────────────────────
def slide_06_horizon_mismatch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Decision 3: Horizon Mismatch \u2014 Don't Change Multiple Hyperparams",
                  "Changed H, n_action_steps, smooth_weight simultaneously \u2192 0% success")
    add_phase_badge(slide, 3, "Feb 17")

    data = [
        ["Run", "Horizon", "n_action_steps", "smooth_weight", "Success", "Issue"],
        ["Baseline", "32", "8", "0.01", "14%", "Working reference"],
        ["Run 2", "48", "16", "0.05", "8%", "Action amplification"],
        ["Run 3", "48", "16", "0.05", "0%", "Action suppression"],
        ["Run 4", "32", "8", "0.01", "14%+", "\u2714 Reverted baseline"],
    ]
    add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.5), data)

    bullets = [
        "Hypothesis was that longer horizon (48) would let the model plan further ahead.",
        "Why it failed: Demos collected with H=32 temporal structure. H=48 creates chunks that "
        "span multiple sub-goals (approach + descent + grasp), confusing the model.",
        "Additional problem: smooth_weight=0.05 (5\u00d7 stronger) over-smoothed action predictions "
        "for the longer horizon \u2014 actions were 30\u201360% of demo standard deviation.",
        "Initial misdiagnosis: Blamed DDPM sampling ('action amplification'). Switching to DDIM "
        "didn't help \u2014 Run 3 got 0%. The real cause was the horizon change.",
        "Lesson: Change ONE hyperparameter at a time. Reverted to baseline config.",
    ]
    add_bullet_slide(slide, Inches(0.4), Inches(4.1), Inches(12.5), Inches(3.2),
                     bullets, font_size=13, color=BLACK)


# ── Slide 7: Bug #3 \u2014 Critical Diffusion Fixes ──────────────────────
def slide_07_diffusion_fixes(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Decision 4: Four Critical Diffusion Bugs \u2014 14% \u2192 84\u201396%",
                  "Feb 22: Research-backed analysis against Chi et al. 2023 reference implementation")
    add_phase_badge(slide, 4, "Feb 22")

    data = [
        ["Bug", "What Was Wrong", "Impact", "Fix"],
        ["#1 tanh on \u03b5",
         "Model output clipped to [-1,1]\nby torch.tanh(). DDPM requires\nunbounded noise predictions.",
         "\u201380%\n(Critical)",
         "Removed tanh().\nModel now predicts\nunbounded \u03b5."],
        ["#2 Normalization",
         "Checkpoint stored act_mean=[0]\nand act_std=[1] instead of\nreal demo statistics.",
         "\u201320%",
         "Saved actual demo\nstatistics in ckpt."],
        ["#3 No EMA",
         "No exponential moving average.\nWeights oscillated during training.",
         "\u201310%",
         "Added EMA with\ndecay=0.999."],
        ["#4 Beta schedule",
         "beta_end=0.02 gave \u03b1\u0304_T=0.36.\nSamples start from structured\nnoise, not pure Gaussian.",
         "\u20135%",
         "Set beta_end=0.1\n\u03b1\u0304_T \u2248 0."],
    ]
    tbl = add_table(slide, Inches(0.3), Inches(1.3), Inches(12.7), Inches(4.0), data)

    add_textbox(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.5),
                "After all four fixes: 84\u201396% task success (up from 14% before fixes, 0% before "
                "execute_steps fix). This is the checkpoint used for all subsequent experiments:\n"
                "runs/diffusion_20260222_195530/ckpt_ep100.pt (101 MB, 8.7M parameters, 100 epochs).",
                font_size=13, bold=True, color=DARK_GREEN)


# ── Slide 8: DDIM Schedule Fix ────────────────────────────────────────
def slide_08_ddim_fix(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Decision 5: DDIM Strided Schedule Fix",
                  "Enabling fast 10-step inference instead of 100-step DDPM")
    add_phase_badge(slide, 4, "Feb 17\u201322")

    bullets = [
        "Original DDIM bug: Assumed consecutive timesteps (t_prev = t \u2212 1),\n"
        "which forced all 100 diffusion steps even with DDIM.",
        "",
        "Fix: Added explicit t_prev parameter to p_sample_ddim().\n"
        "Built strided schedule: [99, 93, 88, \u2026, 0] for 10 steps.",
        "",
        "Validation: Compared single-step DDIM vs DDPM at same timestep:\n"
        "  DDIM/DDPM ratio = 1.03 (within numerical precision).\n"
        "  \u2192 DDIM coefficients are correct.",
        "",
        "The 0.45\u20130.67\u00d7 action suppression seen in full rollouts was\n"
        "cumulative divergence over 100 steps, not a per-step bug.",
        "",
        "Result: DDIM 10 steps matches DDPM 100 steps in quality,\n"
        "with 10\u00d7 faster inference. Used for all subsequent experiments.",
    ]
    add_bullet_slide(slide, Inches(0.4), Inches(1.3), Inches(7.0), Inches(6.0),
                     bullets, font_size=14, color=BLACK)

    # Side table: timing comparison
    timing_data = [
        ["Sampler", "Steps", "Quality"],
        ["DDPM", "100", "Baseline"],
        ["DDIM (old)", "100", "Same (bug forced 100)"],
        ["DDIM (fixed)", "10", "Equivalent"],
    ]
    add_table(slide, Inches(8.0), Inches(1.5), Inches(4.8), Inches(2.2), timing_data)

    add_textbox(slide, Inches(8.0), Inches(4.0), Inches(4.8), Inches(2.0),
                "Single-step diagnostic at multiple\n"
                "timesteps (t=90,70,50,30,10):\n\n"
                "DDIM/DDPM x_prev ratio:\n"
                "1.006, 1.014, 1.028, 1.032, 1.093\n"
                "\u2192 All within 10% = correct",
                font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)


# ── Slide 9: Working Policy Results ────────────────────────────────────
def slide_09_working_policy(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Result: Working Diffusion Policy (84\u201396% Success)",
                  "Checkpoint: diffusion_20260222_195530/ckpt_ep100.pt  |  8.7M params  |  UNet Conv1d")
    add_phase_badge(slide, 5, "Feb 22\u201324")

    # Show composites: straight, moderate, strong
    pairs = [
        ("cfg00_left_arc00.mp4", "Straight (arc 00) \u2014 ambiguous"),
        ("cfg00_left_arc12.mp4", "Moderate arc (arc 12) \u2014 somewhat legible"),
        ("cfg00_left_arc19.mp4", "Strong arc (arc 19) \u2014 highly legible"),
    ]
    for i, (vname, label) in enumerate(pairs):
        y = 1.3 + i * 1.9
        try_add_video_composite(slide, vname, Inches(0.3), Inches(y),
                                Inches(9.0), Inches(1.4), 5)
        add_textbox(slide, Inches(9.5), Inches(y + 0.3), Inches(3.5), Inches(0.8),
                    label, font_size=12, bold=True, color=DARK_BLUE)

    add_textbox(slide, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4),
                "Policy architecture: obs_dim=22, act_dim=5, horizon=32, execute 8 steps, "
                "256-dim hidden, 6 residual blocks, DDIM 10 steps at inference.",
                font_size=11, color=GRAY)

    # Config table on right
    cfg_data = [
        ["Param", "Value"],
        ["Success rate", "84\u201396%"],
        ["Avg steps", "343"],
        ["Epochs", "100"],
        ["Training time", "~1.5 hours"],
    ]
    add_table(slide, Inches(9.5), Inches(1.3), Inches(3.5), Inches(2.5), cfg_data)


# ── Slide 10: Online VLM Steering ──────────────────────────────────────
def slide_10_vlm_steering(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Decision 6: Online VLM Best-of-K Steering",
                  "Inference-time reranking of trajectory candidates using Gemini VLM")
    add_phase_badge(slide, 5, "Feb 24\u201326")

    # Try to add process diagram
    try_add_figure(slide, "vlm_fig3_process_diagram.png",
                   Inches(0.3), Inches(1.3), Inches(6.5), Inches(3.5))

    bullets = [
        "Method: Sample K candidates from diffusion policy with\n"
        "different random seeds \u2192 simulate each for 5 seconds \u2192\n"
        "capture 6 annotated frames \u2192 score with Gemini VLM \u2192\n"
        "execute trajectory with highest legibility score.",
        "",
        "Why this approach? VLM API is non-differentiable, so\n"
        "gradient-based guidance is impossible. Post-hoc reranking\n"
        "works with any scoring function and guarantees in-distribution\n"
        "samples (all candidates came from the trained policy).",
        "",
        "Initial result (K=3, 5 episodes):\n"
        "  \u2022 100% task success (5/5)\n"
        "  \u2022 Mean VLM legibility: 0.750\n"
        "  \u2022 VLM latency: ~5.2 seconds per call",
    ]
    add_bullet_slide(slide, Inches(7.0), Inches(1.3), Inches(6.0), Inches(6.0),
                     bullets, font_size=12, color=BLACK)


# ── Slide 11: VLM Evaluator V1 Failure ──────────────────────────────
def slide_11_vlm_v1_failure(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Decision 7: VLM Evaluator V1 \u2014 45% Accuracy (Random Chance)",
                  "Five compounding bugs made the trajectory evaluator useless")
    add_phase_badge(slide, 6, "Feb 26")

    # V1 vs V2 figure
    try_add_figure(slide, "fig1_v1_vs_v2_accuracy.png",
                   Inches(0.3), Inches(1.3), Inches(6.0), Inches(3.5))

    bugs = [
        "V1 Evaluator \u2014 5 compounding issues:",
        "",
        "1. Camera yaw=135\u00b0 reverses left/right in image",
        "   World LEFT block appears on image RIGHT side.",
        "   VLM had no perspective info \u2192 random mapping.",
        "",
        "2. Prompt was 26,000 characters (XML theory dump)",
        "   Overwhelmed the model with abstract framework.",
        "",
        "3. Temperature = 0.7 (too high for perception)",
        "   Added random variance to deterministic judgments.",
        "",
        "4. No visual anchoring \u2014 raw video only",
        "   VLM had to simultaneously locate blocks +",
        "   track trajectory + determine direction.",
        "",
        "5. Identical block appearance \u2014 both red cubes",
        "   14cm apart, indistinguishable by color.",
        "",
        "Result: 18/40 = 45% \u2014 exactly random chance.",
    ]
    add_bullet_slide(slide, Inches(6.5), Inches(1.3), Inches(6.5), Inches(6.0),
                     bugs, font_size=11, color=BLACK)


# ── Slide 12: VLM Evaluator V2 Fix ─────────────────────────────────
def slide_12_vlm_v2_fix(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Decision 8: VLM V2 Fix \u2014 45% \u2192 97.5% Accuracy",
                  "Five targeted fixes resolved all issues")
    add_phase_badge(slide, 7, "Feb 28")

    # Annotated frame figure
    try_add_figure(slide, "fig6_annotated_frame.png",
                   Inches(0.3), Inches(1.3), Inches(5.5), Inches(3.8))

    fixes = [
        "V2 Fixes (each targeting one V1 bug):",
        "",
        "1. Annotated reference frame with colored markers",
        "   Blue circle = Block A, Green circle = Block B",
        "   VLM sees which block is where in the image.",
        "",
        "2. Prompt reduced to 1,180 characters",
        "   Concise: describe scene, reference image, ask.",
        "",
        "3. Temperature lowered to 0.1",
        "   Near-deterministic perceptual judgments.",
        "",
        "4. Image-space reasoning + coordinate mapping",
        "   VLM answers 'Block A or B?' (image space).",
        "   Eval harness maps to world space using known",
        "   camera calibration: image_left = world_right.",
        "",
        "5. Multi-modal input: annotated frame + video + prompt",
        "",
        "Result: 39/40 = 97.5% accuracy",
        "  LEFT:  19/20 (95%)  |  RIGHT: 20/20 (100%)",
    ]
    add_bullet_slide(slide, Inches(6.0), Inches(1.3), Inches(7.0), Inches(6.0),
                     fixes, font_size=11, color=BLACK)

    # V2 eval IO figure
    try_add_figure(slide, "fig4_v2_eval_io.png",
                   Inches(0.3), Inches(5.3), Inches(5.5), Inches(1.8))

    # Confusion matrix table
    cm_data = [
        ["", "Pred LEFT", "Pred RIGHT"],
        ["GT LEFT", "19", "1"],
        ["GT RIGHT", "0", "20"],
    ]
    add_table(slide, Inches(6.2), Inches(5.5), Inches(3.5), Inches(1.5), cm_data)

    add_textbox(slide, Inches(10.0), Inches(5.5), Inches(3.0), Inches(1.5),
                "Key insight: The VLM's\nperception was never the\nbottleneck \u2014 the failure\nwas entirely a grounding\nproblem.",
                font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)


# ── Slide 13: VLM Debug Marathon ─────────────────────────────────────
def slide_13_vlm_debug_marathon(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "The VLM Debug Marathon \u2014 23 Runs in One Day",
                  "Feb 28: Systematic testing to understand why overhead camera struggles")
    add_phase_badge(slide, 7, "Feb 28")

    bullets = [
        "After V2 evaluator reached 97.5%, we tested how well it could differentiate\n"
        "between arc magnitudes (not just left vs right).",
        "",
        "23 debug runs (vlm_debug_20260228_*) explored:",
        "  \u2022 Different camera viewpoints (overhead, angled, front)",
        "  \u2022 Different arc class pairs (arc04 vs arc10, arc04 vs arc15, arc10 vs arc15)",
        "  \u2022 Different temporal sampling (t=1s, 2s, 3s, 4s, 5s, 6s)",
        "  \u2022 Pairwise perception tests with side-by-side image comparisons",
        "",
        "Key finding: The overhead camera (yaw=135\u00b0) compresses lateral displacement.",
        "Arc 04 vs arc 15 look very similar from above. The VLM can distinguish",
        "LEFT vs RIGHT reliably, but cannot rank arc magnitudes precisely.",
        "",
        "This motivated the shift from 'VLM ranks arcs' to 'VLM confirms direction,\n"
        "geometric metric ranks arc magnitude' \u2014 the hybrid approach used in the",
        "paired replanning experiments.",
    ]
    add_bullet_slide(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(6.0),
                     bullets, font_size=13, color=BLACK)


# ── Slide 14: Paired Replanning Results ──────────────────────────────
def slide_14_paired_replanning(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Result: Paired Replanning \u2014 +69.4% Arc Improvement",
                  "4 paired episodes: same seed, baseline vs VLM-steered, 100% success both sides")
    add_phase_badge(slide, 8, "Feb 28")

    data = [
        ["Pair", "Block", "Baseline Arc", "Guided Arc", "Improvement", "Legibility"],
        ["1", "RIGHT", "0.631 m", "0.635 m", "+0.6%", "1.000"],
        ["2", "RIGHT", "0.033 m", "0.750 m", "+2204%", "0.990"],
        ["3", "LEFT", "0.403 m", "0.891 m", "+121%", "0.850"],
        ["4", "LEFT", "0.690 m", "0.701 m", "+1.5%", "0.950"],
        ["Mean", "\u2014", "0.439 m", "0.744 m", "+69.4%", "0.948"],
    ]
    tbl = add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.3), data)

    bullets = [
        "Same-block accuracy: 100% \u2014 VLM-guided always matched baseline's target block.",
        "Arc-15-19 selection: 100% \u2014 all guided trajectories selected legible arcs.",
        "Task success: 100% for both baseline and guided (all 8 episodes succeeded).",
        "Pair 2 shows the biggest impact: baseline was nearly straight (0.033m),\n"
        "guided found a strongly curved alternative (0.750m) to the same block.",
        "Selection method: arc15_legible \u2014 among candidates with VLM legibility \u2265 0.70,\n"
        "pick the one with highest geometric arc magnitude.",
    ]
    add_bullet_slide(slide, Inches(0.4), Inches(4.9), Inches(12.5), Inches(2.5),
                     bullets, font_size=12, color=BLACK)


# ── Slide 15: 50-Episode Steering Experiment ─────────────────────────
def slide_15_full_steering_experiment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Result: 50-Episode VLM Steering Experiment",
                  "Baseline (no steering) vs VLM-steered (K=10, proxy scorer)")
    add_phase_badge(slide, 8, "Mar 1\u20139")

    # Try arc distribution figure
    try_add_figure(slide, "vlm_fig1_arc_distribution.png",
                   Inches(0.3), Inches(1.3), Inches(6.2), Inches(3.5))

    data = [
        ["Metric", "Baseline", "VLM-Steered", "Change"],
        ["Episodes", "50", "50", "\u2014"],
        ["Task Success", "94% (47/50)", "98% (proxy)", "+4%"],
        ["Mean Arc", "0.080 m", "0.114 m", "+42.5%"],
        ["Arc 00\u201304 (straight)", "58%", "20%", "\u221238 pp"],
        ["Arc 15\u201319 (strong)", "6%", "34%", "+28 pp"],
        ["VLM Legibility", "0.604", "0.946", "+56.6%"],
        ["Steering Acc.", "\u2014", "100%", "\u2014"],
    ]
    add_table(slide, Inches(6.8), Inches(1.3), Inches(6.2), Inches(4.2), data)

    add_textbox(slide, Inches(0.4), Inches(5.2), Inches(12.5), Inches(2.0),
                "Without steering, 58% of trajectories are straight (arc 00\u201304) and ambiguous. "
                "VLM steering shifts the distribution: only 20% straight, 34% strongly curved.\n"
                "Proxy scorer achieves comparable steering at 98% success vs 70% for pure VLM "
                "(because VLM sometimes selects visually legible but task-risky candidates).",
                font_size=12, color=BLACK)


# ── Slide 16: Success-Legibility Trade-off ──────────────────────────
def slide_16_tradeoff(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Decision 9: The Success\u2013Legibility Trade-off",
                  "Pure VLM steering drops success to 70% \u2014 this motivated LegDiff")
    add_phase_badge(slide, 8, "Mar 9")

    try_add_figure(slide, "vlm_fig4_success_legibility.png",
                   Inches(0.3), Inches(1.3), Inches(6.5), Inches(4.5))

    bullets = [
        "The core problem: Best-of-K VLM steering improves\n"
        "legibility by +57% but drops success from 94% to 70%.",
        "",
        "Why? The VLM selects candidates that look legible but\n"
        "may be on the edge of the policy's capability. Strongly\n"
        "curved arcs are harder to execute successfully.",
        "",
        "Three solutions explored:",
        "  1. Proxy scorer (8 geometric features) \u2192 98% success\n"
        "     but requires environment simulation for candidates",
        "  2. Best-of-16 sampling \u2192 L_post=0.852 (+10.1%)\n"
        "     but modest legibility gain, expensive",
        "  3. LegDiff with CFG (chosen) \u2192 100% success +\n"
        "     L_early=0.934 \u2014 no trade-off, trains a new model",
        "",
        "Decision: Pursue LegDiff as the principled solution.\n"
        "It learns goal-conditioned legibility during training,\n"
        "so inference-time steering doesn't hurt success.",
    ]
    add_bullet_slide(slide, Inches(7.0), Inches(1.3), Inches(6.0), Inches(6.0),
                     bullets, font_size=12, color=BLACK)


# ── Slide 17: Best-of-K Study ────────────────────────────────────────
def slide_17_best_of_k(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Best-of-K Analysis: How Much Diversity Exists?",
                  "Sampling K candidates from the same observation")
    add_phase_badge(slide, 8, "Mar 9")

    try_add_figure(slide, "vlm_fig5_candidate_selection.png",
                   Inches(0.3), Inches(1.3), Inches(6.5), Inches(4.0))

    # Best-of-N results
    data = [
        ["K", "L_early", "L_posterior", "Relative Gain"],
        ["1 (baseline)", "0.732", "0.789", "\u2014"],
        ["2", "0.756", "0.810", "+2.7%"],
        ["4", "0.778", "0.831", "+5.3%"],
        ["8", "0.793", "0.844", "+7.0%"],
        ["16", "0.806", "0.852", "+10.1%"],
    ]
    add_table(slide, Inches(7.2), Inches(1.3), Inches(5.8), Inches(3.0), data)

    bullets = [
        "Diminishing returns: Going from K=1 to K=16 only gains 10.1%.",
        "The policy's natural diversity is limited \u2014 most candidates cluster",
        "around similar arc magnitudes for the same observation.",
        "",
        "Arc diversity analysis (100 rollouts): mean_entropy=0.284,",
        "87% failure rate, only 13% picked either block reliably.",
        "The policy collapsed to producing similar trajectories.",
        "",
        "This showed that post-hoc reranking has a ceiling.",
        "For larger legibility gains, we need to modify the",
        "generation process itself \u2192 LegDiff.",
    ]
    add_bullet_slide(slide, Inches(7.2), Inches(4.5), Inches(5.8), Inches(3.0),
                     bullets, font_size=11, color=BLACK)


# ── Slide 18: LegDiff Architecture ──────────────────────────────────
def slide_18_legdiff_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Decision 10: LegDiff \u2014 Classifier-Free Guidance for Legibility",
                  "Goal-conditioned diffusion policy with null-goal dropout (p_uncond=0.15)")
    add_phase_badge(slide, 10, "Mar 30\u201331")

    bullets_left = [
        "Architecture: GoalCondDiffusionPolicy",
        "  \u2022 Same UNet Conv1d backbone as baseline",
        "  \u2022 Added: Goal embedding (3 tokens: left/right/null)",
        "  \u2022 During training: 15% of batches use null goal (dropout)",
        "  \u2022 At inference: CFG formula steers toward specified goal:",
        "    \u03b5\u0303 = \u03b5_uncond + w \u00b7 (\u03b5_cond \u2212 \u03b5_uncond)",
        "    where w = 3.0 (guidance scale)",
        "",
        "Training details:",
        "  \u2022 150 epochs, batch_size=256, lr=1e-4",
        "  \u2022 Mirror augmentation (doubles dataset, flips goal labels)",
        "  \u2022 Weighted sampling (priority on early timesteps t < 10%)",
        "  \u2022 Loss: 0.047 \u2192 0.0018 (96.3% reduction over 100 epochs)",
        "",
        "Key reference: Bronars et al. (RA-L 2024)",
        "  'Legibility Diffuser: Offline Imitation for",
        "   Intent-Expressive Motion'",
    ]
    add_bullet_slide(slide, Inches(0.4), Inches(1.3), Inches(6.5), Inches(6.0),
                     bullets_left, font_size=12, color=BLACK)

    # Architecture diagram as text
    arch_bullets = [
        "Model Structure:",
        "  Time: Sinusoidal(128) \u2192 Linear(256)",
        "  Obs:  Linear(22 \u2192 256)",
        "  Goal: Embedding(3 \u2192 256)",
        "  Input: Conv1d(5 \u2192 256)",
        "  Encoder: ResBlock(256\u2192256) \u2192 ResBlock(256\u2192512)",
        "  Bottleneck: ResBlock(512\u2192512)",
        "  Decoder: ResBlock(1024\u2192512) \u2192 ResBlock(512\u2192256)",
        "  Output: Conv1d(256 \u2192 5)",
        "",
        "Checkpoint: 34 MB (vs 101 MB baseline)",
        "Smaller because no EMA weights saved.",
    ]
    add_bullet_slide(slide, Inches(7.2), Inches(1.3), Inches(5.8), Inches(3.2),
                     arch_bullets, font_size=11, color=GRAY)

    # Training loss mini-table
    loss_data = [
        ["Epoch", "Loss"],
        ["1", "0.0469"],
        ["10", "0.0075"],
        ["50", "0.0027"],
        ["100", "0.0018"],
    ]
    add_table(slide, Inches(8.5), Inches(4.8), Inches(3.5), Inches(2.2), loss_data)


# ── Slide 19: LegDiff Results ────────────────────────────────────────
def slide_19_legdiff_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Result: LegDiff \u2014 100% Success + High Legibility, No Trade-off",
                  "20 episodes (CFG only) + 10 episodes (CFG + VLM K=3)")
    add_phase_badge(slide, 10, "Mar 31")

    data = [
        ["Method", "Episodes", "Success", "L_early (mean)", "L_early (std)", "Goal Match"],
        ["Baseline", "20", "100%", "0.919", "\u00b10.029", "\u2014"],
        ["LegDiff CFG w=3.0", "20", "100%", "0.935", "\u00b10.020", "100%"],
        ["LegDiff + VLM K=3", "10", "100%", "0.955", "\u00b10.013", "100%"],
    ]
    tbl = add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.2), data)

    # Highlight the LegDiff rows
    table = tbl.table
    for c in range(6):
        table.cell(2, c).fill.solid()
        table.cell(2, c).fill.fore_color.rgb = LIGHT_GREEN
        table.cell(3, c).fill.solid()
        table.cell(3, c).fill.fore_color.rgb = LIGHT_GREEN

    bullets = [
        "LegDiff achieves the thesis goal: increased legibility without sacrificing success.",
        "",
        "Comparison to previous best results:",
        "  \u2022 Online VLM Best-of-K: 70% success, 0.946 VLM legibility \u2014 big success drop",
        "  \u2022 Proxy scorer: 98% success, 100% steering accuracy \u2014 good but needs simulator",
        "  \u2022 LegDiff CFG: 100% success, L_early=0.935 \u2014 no trade-off, no simulator needed",
        "  \u2022 LegDiff + VLM: 100% success, L_early=0.955 \u2014 best of all methods",
        "",
        "Why LegDiff resolves the trade-off:",
        "  The policy was TRAINED to produce goal-directed legible trajectories.",
        "  CFG amplifies the goal signal during denoising, pushing the trajectory",
        "  toward one that clearly reveals the intended goal. Because this is baked",
        "  into the sampling process, there is no separate selection step that",
        "  might pick task-risky candidates.",
    ]
    add_bullet_slide(slide, Inches(0.4), Inches(3.8), Inches(12.5), Inches(3.5),
                     bullets, font_size=12, color=BLACK)


# ── Slide 20: Complete Results Summary Table ─────────────────────────
def slide_20_complete_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Complete Results: Every Method Compared")

    data = [
        ["#", "Method", "Date", "Success", "Legibility", "Key Finding"],
        ["1", "Base diffusion", "Feb 22", "84\u201396%", "\u2014", "Working policy after 4 bug fixes"],
        ["2", "VLM K=3 (online)", "Feb 26", "100% (n=5)", "0.750", "Proof of concept, expensive"],
        ["3", "VLM v1 evaluator", "Feb 26", "45% acc.", "\u2014", "Random chance \u2014 5 bugs"],
        ["4", "VLM v2 evaluator", "Feb 28", "97.5% acc.", "0.919", "Grounding fix resolved it"],
        ["5", "Paired replan (4)", "Feb 28", "100%", "0.948", "+69.4% arc, same block"],
        ["6", "Proxy scorer (50)", "Mar 1", "98%", "100% steer", "Geom features, <1ms"],
        ["7", "Best-of-16", "Mar 9", "\u2014", "0.852", "+10.1% \u2014 ceiling on diversity"],
        ["8", "LegDiff CFG (20)", "Mar 31", "100%", "0.935", "No trade-off!"],
        ["9", "LegDiff+VLM (10)", "Mar 31", "100%", "0.955", "Best overall"],
    ]
    tbl = add_table(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.0), data)

    # Highlight LegDiff rows
    table = tbl.table
    for c in range(6):
        table.cell(8, c).fill.solid()
        table.cell(8, c).fill.fore_color.rgb = LIGHT_GREEN
        table.cell(9, c).fill.solid()
        table.cell(9, c).fill.fore_color.rgb = LIGHT_GREEN

    add_textbox(slide, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.6),
                "Evolution of the approach: Online VLM \u2192 VLM evaluator fix \u2192 paired replanning "
                "\u2192 proxy scorer \u2192 diversity ceiling hit \u2192 LegDiff resolves the trade-off.",
                font_size=12, bold=True, color=DARK_GREEN)

    try_add_figure(slide, "fig5_final_results.png",
                   Inches(0.3), Inches(1.2), Inches(0), Inches(0))  # skip this if overlaps


# ── Slide 21: Decision Flow ─────────────────────────────────────────
def slide_21_decision_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Decision Flow: Why Each Step Led to the Next")

    flow = [
        "1. Architecture search (Feb 8\u201313): 7 failed runs \u2192 UNet Conv1d works",
        "   WHY: Small MLPs can't represent 400-demo multimodal distribution",
        "",
        "2. execute_steps fix (Feb 16): 0% \u2192 50% success",
        "   WHY: Training temporal structure must match eval",
        "",
        "3. Horizon revert (Feb 17): 0\u20138% \u2192 14% success",
        "   WHY: Changed 3 hyperparams at once, couldn't isolate cause",
        "",
        "4. Diffusion bug fixes (Feb 22): 14% \u2192 84\u201396% success",
        "   WHY: tanh on noise, wrong norm, no EMA, bad beta schedule",
        "",
        "5. VLM steering (Feb 24): Works but expensive (\u223c5s/call)",
        "   WHY: VLM is non-differentiable, so reranking is the only option",
        "",
        "6. VLM evaluator fix (Feb 28): 45% \u2192 97.5% accuracy",
        "   WHY: Camera yaw=135\u00b0 reverses image, needed visual anchoring",
        "",
        "7. Paired replanning (Feb 28): +69.4% arc, 0.948 legibility",
        "   WHY: Validated steering preserves task while improving legibility",
        "",
        "8. Diversity ceiling (Mar 9): Best-of-16 gains only +10.1%",
        "   WHY: Policy naturally clusters around similar arcs",
        "",
        "9. LegDiff CFG (Mar 30\u201331): 100% success, 0.935 legibility",
        "   WHY: Trains goal-awareness into the model \u2192 no reranking needed",
    ]
    add_bullet_slide(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(6.2),
                     flow, font_size=11, color=BLACK, spacing=Pt(2))


# ── Slide 22: Timeline Visualization ─────────────────────────────────
def slide_22_timeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Project Timeline: Feb 8 \u2013 Mar 31, 2026")

    # Timeline as a table
    data = [
        ["Phase", "Dates", "Work Done", "Key Outcome"],
        ["1", "Feb 8\u201313", "Demo collection + architecture search", "400 demos, UNet Conv1d selected"],
        ["2", "Feb 13\u201316", "Action collapse debugging", "execute_steps=8 fix, 0\u219250%"],
        ["3", "Feb 16\u201317", "Horizon mismatch", "Reverted H=32, n=8"],
        ["4", "Feb 17\u201322", "DDIM + diffusion fixes", "4 bugs fixed, 84\u201396% success"],
        ["5", "Feb 22\u201324", "VLM steering implementation", "Best-of-K, Gemini integration"],
        ["6", "Feb 24\u201326", "VLM evaluator V1", "45% accuracy, 5 bugs found"],
        ["7", "Feb 26\u201328", "VLM V2 + debug marathon", "97.5% accuracy, 23 debug runs"],
        ["8", "Feb 28\u2013Mar 9", "Paired replanning + proxy scorer", "+69.4% arc, 98% success"],
        ["9", "Mar 9", "Best-of-K analysis", "Diversity ceiling at +10.1%"],
        ["10", "Mar 30\u201331", "LegDiff CFG training + eval", "100% success, L=0.935"],
    ]
    tbl = add_table(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5), data)

    # Color-code phases
    table = tbl.table
    # Debug phases in light red
    for phase in [2, 3, 4, 7]:
        for c in range(4):
            table.cell(phase, c).fill.solid()
            table.cell(phase, c).fill.fore_color.rgb = LIGHT_RED
    # Result phases in light green
    for phase in [5, 8, 9, 10]:
        for c in range(4):
            table.cell(phase, c).fill.solid()
            table.cell(phase, c).fill.fore_color.rgb = LIGHT_GREEN

    add_textbox(slide, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.4),
                "Red = debugging phases  |  Green = implementation & results phases  |  "
                "Total: 12 training runs, 30+ VLM debug runs, 67 experiment directories",
                font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)


# ── Slide 23: Limitations ────────────────────────────────────────────
def slide_23_limitations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Honest Assessment: Limitations & Open Questions")

    lim_bullets = [
        "1. Small evaluation scale: LegDiff tested on 20+10 episodes.\n"
        "   Larger-scale evaluation (100+ episodes) needed for confidence.",
        "",
        "2. Single environment: TwoBlockPick with two identical blocks.\n"
        "   Generalisation to different objects, more goals, longer\n"
        "   horizons, or real robot hardware is not yet shown.",
        "",
        "3. Legibility metric is geometric (arc displacement), not\n"
        "   human-validated. Human studies needed to confirm that\n"
        "   arced trajectories are actually perceived as more legible.",
        "",
        "4. VLM calls are expensive (\u223c5s each). The LegDiff approach\n"
        "   avoids runtime VLM, but training requires goal labels which\n"
        "   were derived from the demo structure (left/right block).",
        "",
        "5. No real-robot experiments. All results are in PyBullet\n"
        "   simulation. Sim-to-real transfer is an open question.",
        "",
        "6. CFG scale w=3.0 was chosen based on limited tuning.\n"
        "   A systematic sweep over w could improve results further.",
    ]
    add_bullet_slide(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(6.0),
                     lim_bullets, font_size=13, color=BLACK)


# ── Slide 24: What's Next ────────────────────────────────────────────
def slide_24_whats_next(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "What Further Work Is Expected",
                  "Concrete next steps toward thesis completion")

    next_bullets = [
        "1. Large-scale LegDiff evaluation (100+ episodes)",
        "   Establish statistically significant results with confidence intervals.",
        "",
        "2. CFG scale sweep (w = 1.0, 2.0, 3.0, 5.0, 7.0)",
        "   Characterise the success\u2013legibility Pareto frontier for LegDiff.",
        "",
        "3. Human study",
        "   Show trajectory videos to participants; measure time-to-prediction",
        "   of the robot's goal. Validate that geometric arc \u2248 perceptual legibility.",
        "",
        "4. Additional environments / robot tasks",
        "   Test on 3+ objects, different workspaces, or longer-horizon tasks",
        "   to assess generality of the approach.",
        "",
        "5. Sim-to-real transfer",
        "   Deploy on real Franka Panda (if available) or comparable hardware.\n"
        "   Verify VLM + LegDiff approach works outside simulation.",
        "",
        "6. Thesis writing",
        "   Draft chapters: Background, Methodology, Experiments, Analysis.",
        "   All experimental infrastructure and results exist to support writing.",
    ]
    add_bullet_slide(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(6.0),
                     next_bullets, font_size=13, color=BLACK)


# ── Slide 25: Summary ────────────────────────────────────────────────
def slide_25_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BLUE)

    add_textbox(slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.7),
                "Summary",
                font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(1.3), Inches(5.333), Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

    takeaways = [
        "1.  Built a working diffusion policy for TwoBlockPick (84\u201396% success)\n"
        "     after fixing 4 critical bugs in the diffusion implementation.",
        "",
        "2.  Implemented VLM-based trajectory evaluation and steering,\n"
        "     fixing the evaluator from 45% \u2192 97.5% accuracy along the way.",
        "",
        "3.  Demonstrated that Best-of-K VLM steering improves legibility (+57%)\n"
        "     but has a success\u2013legibility trade-off (100% \u2192 70% success).",
        "",
        "4.  Identified a diversity ceiling (best-of-16 gains only +10.1%),\n"
        "     motivating a fundamentally different approach.",
        "",
        "5.  LegDiff (Classifier-Free Guidance) resolves the trade-off:\n"
        "     100% success + L_early = 0.935 (CFG) / 0.955 (CFG + VLM).\n"
        "     No runtime VLM needed, no simulator needed, no success drop.",
        "",
        "6.  Complete experimental pipeline built: 67 run directories,\n"
        "     13 JSON result files, 16 publication-ready figures, CLI tooling.",
    ]
    add_bullet_slide(slide, Inches(0.6), Inches(1.7), Inches(12.0), Inches(5.5),
                     takeaways, font_size=14, color=WHITE, spacing=Pt(3))


# ═════════════════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ═════════════════════════════════════════════════════════════════════════

def main():
    print("\n" + "=" * 65)
    print("  Building Thesis Progress Presentation v2")
    print("=" * 65 + "\n")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        ("Title", slide_01_title),
        ("Research Context & Goal", slide_02_research_context),
        ("Task, Environment & Data", slide_03_task_and_data),
        ("Architecture Search (12 runs)", slide_04_architecture_search),
        ("Bug: execute_steps", slide_05_execute_steps_bug),
        ("Bug: Horizon Mismatch", slide_06_horizon_mismatch),
        ("Bug: 4 Diffusion Fixes", slide_07_diffusion_fixes),
        ("DDIM Schedule Fix", slide_08_ddim_fix),
        ("Working Policy (84-96%)", slide_09_working_policy),
        ("Online VLM Steering", slide_10_vlm_steering),
        ("VLM V1 Failure (45%)", slide_11_vlm_v1_failure),
        ("VLM V2 Fix (97.5%)", slide_12_vlm_v2_fix),
        ("VLM Debug Marathon", slide_13_vlm_debug_marathon),
        ("Paired Replanning (+69.4%)", slide_14_paired_replanning),
        ("50-Episode Experiment", slide_15_full_steering_experiment),
        ("Success-Legibility Trade-off", slide_16_tradeoff),
        ("Best-of-K Analysis", slide_17_best_of_k),
        ("LegDiff Architecture", slide_18_legdiff_architecture),
        ("LegDiff Results", slide_19_legdiff_results),
        ("Complete Results Table", slide_20_complete_results),
        ("Decision Flow", slide_21_decision_flow),
        ("Timeline", slide_22_timeline),
        ("Limitations", slide_23_limitations),
        ("What's Next", slide_24_whats_next),
        ("Summary", slide_25_summary),
    ]

    for name, builder in builders:
        try:
            builder(prs)
            print(f"  \u2713 Slide: {name}")
        except Exception as e:
            print(f"  \u2717 Slide: {name} \u2014 ERROR: {e}")

    prs.save(str(PPTX_PATH))
    print(f"\n  Saved \u2192 {PPTX_PATH}")
    print(f"  Total slides: {len(prs.slides)}")
    print("=" * 65 + "\n")


if __name__ == "__main__":
    main()
