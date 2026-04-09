#!/usr/bin/env python
r"""
Build clean method slides for thesis presentation.
Output: presentation/method_slides.pptx
Usage: .venv\Scripts\python.exe scripts/build_method_slides.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT    = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / "presentation"
OUT_DIR.mkdir(exist_ok=True)
PPTX_PATH = OUT_DIR / "method_slides.pptx"

NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
LBLUE  = RGBColor(0xDB, 0xEA, 0xFE)
TEAL   = RGBColor(0x0D, 0x94, 0x88)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
LGREEN = RGBColor(0xD1, 0xFA, 0xE5)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
LPURP  = RGBColor(0xED, 0xE9, 0xFE)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
LORAN  = RGBColor(0xFF, 0xED, 0xD5)
GOLD   = RGBColor(0xF5, 0x9E, 0x0B)
LGOLD  = RGBColor(0xFF, 0xFB, 0xEB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x11, 0x18, 0x27)
LGRAY  = RGBColor(0xF1, 0xF5, 0xF9)
MGRAY  = RGBColor(0x94, 0xA3, 0xB8)
RED    = RGBColor(0xDC, 0x26, 0x26)
LRED   = RGBColor(0xFE, 0xE2, 0xE2)


def bg(slide, color):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, line=None, lw=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, size=16, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Calibri"
    return box


def multiline(slide, l, t, w, h, lines, size=14, color=BLACK, gap=Pt(5)):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = gap
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.name = "Calibri"


def header(slide, title, sub=None):
    rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
    txt(slide, Inches(0.4), Inches(0.07), Inches(12.5), Inches(0.54),
        title, size=25, bold=True, color=WHITE)
    if sub:
        txt(slide, Inches(0.4), Inches(0.61), Inches(12.5), Inches(0.34),
            sub, size=12, italic=True, color=LBLUE)


def pill(slide, l, t, text, fill=BLUE):
    rect(slide, l, t, Inches(1.55), Inches(0.28), fill)
    txt(slide, l + Inches(0.07), t + Inches(0.02), Inches(1.5), Inches(0.26),
        text, size=10, bold=True, color=WHITE)


def fcard(slide, l, t, w, h, label, formula, fill=LGRAY, lc=BLUE):
    rect(slide, l, t, w, h, fill, lc, Pt(1.5))
    if label:
        txt(slide, l+Inches(0.12), t+Inches(0.06), w-Inches(0.2), Inches(0.22),
            label, size=9, bold=True, color=lc)
    yf = t+Inches(0.28) if label else t+Inches(0.08)
    hf = h - Inches(0.34) if label else h - Inches(0.14)
    txt(slide, l+Inches(0.12), yf, w-Inches(0.2), hf,
        formula, size=14, color=BLACK, wrap=True)


def stepbox(slide, l, t, w, h, num, title, note, fill=LBLUE):
    rect(slide, l, t, w, h, fill, BLUE, Pt(1))
    rect(slide, l+Inches(0.1), t+Inches(0.1), Inches(0.28), Inches(0.28), BLUE)
    txt(slide, l+Inches(0.1), t+Inches(0.08), Inches(0.28), Inches(0.28),
        str(num), size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, l+Inches(0.45), t+Inches(0.06), w-Inches(0.54), Inches(0.3),
        title, size=12, bold=True, color=NAVY)
    txt(slide, l+Inches(0.45), t+Inches(0.35), w-Inches(0.54), h-Inches(0.41),
        note, size=10, color=BLACK, wrap=True)


# ── SLIDES ─────────────────────────────────────────────────────────────

def s0_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    txt(slide, Inches(1.0), Inches(0.7), Inches(11.3), Inches(1.1),
        "Multimodal Diffusion Policy", size=36, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, Inches(1.5), Inches(1.85), Inches(10.3), Inches(0.45),
        "Teaching a robot arm to move expressively using diffusion models",
        size=16, italic=True, color=LBLUE, align=PP_ALIGN.CENTER)

    steps = [
        ("1  Demo Data", BLUE),
        ("2  DDPM Train", TEAL),
        ("3  CFG Trick",  PURPLE),
        ("4  U-Net",      GREEN),
        ("5  DDIM",       ORANGE),
        ("6  VLM Judge",  RED),
        ("7  Metrics",    GOLD),
    ]
    bw = Inches(1.6); bh = Inches(0.75); gap = Inches(0.14)
    total = len(steps)*bw + (len(steps)-1)*gap
    sx = (Inches(13.333) - total) / 2
    for i, (label, col) in enumerate(steps):
        x = sx + i*(bw+gap)
        rect(slide, x, Inches(3.1), bw, bh, col)
        txt(slide, x, Inches(3.1)+Inches(0.18), bw, bh,
            label, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(steps)-1:
            rect(slide, x+bw, Inches(3.1)+bh//2-Inches(0.03), gap, Inches(0.06), WHITE)

    txt(slide, Inches(0.5), Inches(4.15), Inches(12.3), Inches(0.35),
        "train_cfg.py  (slides 2-4)                             eval_cfg_vlm.py  (slides 5-7)",
        size=12, color=LBLUE, align=PP_ALIGN.CENTER)
    txt(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
        "Franka Panda  .  PyBullet  .  Two-Block Pick Task",
        size=13, italic=True, color=MGRAY, align=PP_ALIGN.CENTER)


def s1_ddpm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    header(slide, "DDPM Training — The Robot Learns by Denoising",
           "Add noise to demonstrations, train the network to undo it")
    pill(slide, Inches(11.9), Inches(0.12), "2 / 7", TEAL)

    txt(slide, Inches(0.4), Inches(1.12), Inches(5.6), Inches(0.3),
        "The Idea  (3 steps)", size=14, bold=True, color=NAVY)
    multiline(slide, Inches(0.4), Inches(1.44), Inches(5.6), Inches(1.1),
        ["1.  Take a clean action chunk  a0",
         "2.  Add random noise  up to T=100 levels",
         "3.  Train the network: given noisy a\u209c + level t, predict the noise"],
        size=13, color=BLACK)

    fcard(slide, Inches(0.4), Inches(2.65), Inches(5.6), Inches(0.75),
          "Noise Schedule   (\u03b1\u0305\u209c = how much signal remains at step t)",
          "\u03b2\u2081\u2026\u03b2\u209c linearly  0.0001 \u2192 0.1       "
          "\u03b1\u0305\u209c = \u03b2\u2081 x \u03b2\u2082 x ... x \u03b2\u209c")

    fcard(slide, Inches(0.4), Inches(3.5), Inches(5.6), Inches(0.75),
          "Forward Process   (add all noise in ONE shot)",
          "a\u209c  =  \u221a\u03b1\u0305\u209c \u00b7 a\u2080  +  \u221a(1 \u2212 \u03b1\u0305\u209c) \u00b7 \u03b5        "
          "\u03b5 ~ N(0, I)")

    fcard(slide, Inches(0.4), Inches(4.35), Inches(5.6), Inches(0.85),
          "Training Loss   (mean squared error on the predicted noise)",
          "L  =  E [ || \u03b5  -  \u03b5\u03b8(a\u209c, t, obs) ||^2 ]"
          "\n\u03b5\u03b8 = network prediction,    obs = robot observation",
          fill=LGREEN, lc=GREEN)

    fcard(slide, Inches(0.4), Inches(5.32), Inches(5.6), Inches(0.65),
          "EMA Shadow Weights   (saved; used at eval time)",
          "\u03b8_ema  <-  0.999 x \u03b8_ema  +  0.001 x \u03b8")

    txt(slide, Inches(6.4), Inches(1.12), Inches(6.6), Inches(0.3),
        "Training Loop  (one batch)", size=14, bold=True, color=NAVY)
    loop = [
        ("1", "Pick random noise level  t",    "t ~ Uniform(0, 100)  each batch"),
        ("2", "Corrupt the action  a0 -> at",  "One-shot formula above"),
        ("3", "Network predicts noise  e_hat",  "e_hat = e_theta(at, t, obs)"),
        ("4", "Loss = MSE(e_hat, e)",           "Backprop + AdamW update"),
        ("5", "Update EMA weights",             "Slow-moving copy of the model"),
    ]
    for i, (n, t_, note) in enumerate(loop):
        stepbox(slide, Inches(6.4), Inches(1.45)+i*Inches(0.95), Inches(6.6), Inches(0.83), n, t_, note)

    txt(slide, Inches(0.4), Inches(6.58), Inches(12.6), Inches(0.28),
        "The network solves just ONE random denoising step per batch - not the full 100-step chain.",
        size=11, italic=True, color=MGRAY)


def s2_cfg(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    header(slide, "Classifier-Free Guidance (CFG)",
           "One model, four behaviors - no separate models needed")
    pill(slide, Inches(11.9), Inches(0.12), "3 / 7", PURPLE)

    txt(slide, Inches(0.4), Inches(1.12), Inches(5.8), Inches(0.3),
        "The Problem", size=14, bold=True, color=NAVY)
    txt(slide, Inches(0.4), Inches(1.44), Inches(5.8), Inches(0.55),
        "We want the robot to behave differently (legibly, safely...) "
        "using ONE trained model. Behavior is encoded in obs dims 22-25.",
        size=13, color=BLACK, wrap=True)

    txt(slide, Inches(0.4), Inches(2.1), Inches(5.8), Inches(0.3),
        "The Training Trick  (CFG dropout, p=0.15)", size=14, bold=True, color=PURPLE)
    fcard(slide, Inches(0.4), Inches(2.42), Inches(5.8), Inches(0.78),
          "Randomly hide the behavior signal during training",
          "obs_tilde = obs               with prob 0.85   (behavior-aware)"
          "\nobs_tilde = obs_uncond    with prob 0.15   (behavior-blind)",
          fill=LPURP, lc=PURPLE)

    txt(slide, Inches(0.4), Inches(3.32), Inches(5.8), Inches(0.3),
        "At Inference  (steer toward a behavior)", size=14, bold=True, color=NAVY)
    fcard(slide, Inches(0.4), Inches(3.64), Inches(5.8), Inches(0.68),
          "Blend the two predictions with guidance scale lambda",
          "e_hat  =  e(obs_uncond)  +  lambda x ( e(obs_cond) - e(obs_uncond) )"
          "\nlambda=0: no steering    lambda=2-4: strong behavior")

    txt(slide, Inches(0.4), Inches(4.45), Inches(2.5), Inches(0.28),
        "lambda values:", size=12, bold=True, color=NAVY)
    multiline(slide, Inches(0.4), Inches(4.75), Inches(5.8), Inches(0.9),
        ["  lambda = 0   ignore behavior, use generic actions",
         "  lambda = 1   standard conditional",
         "  lambda = 2-4   strongly steered behavior"],
        size=12, color=BLACK)

    txt(slide, Inches(6.6), Inches(1.12), Inches(6.4), Inches(0.3),
        "Behavior Conditioning Dims", size=14, bold=True, color=NAVY)

    rows = [
        (LBLUE,  NAVY,   "Dims 0-21",    "Base robot state — always present"),
        (LGREEN, GREEN,  "Dims 22-24",   "Context XYZ (obstacle or waypoint position)"),
        (LPURP,  PURPLE, "Dim 25",       "Behavior mode number"),
        (LGOLD,  GOLD,   "Legibility",   "mode=1,  no context"),
        (LORAN,  ORANGE, "Safety",       "mode=3,  context = obstacle XYZ"),
        (LRED,   RED,    "Grounding",    "mode=4,  context = waypoint XYZ"),
    ]
    for i, (fill, col, name, desc) in enumerate(rows):
        y = Inches(1.46) + i*Inches(0.78)
        rect(slide, Inches(6.6), y, Inches(6.4), Inches(0.68), fill, col, Pt(1))
        txt(slide, Inches(6.75), y+Inches(0.07), Inches(2.1), Inches(0.3),
            name, size=12, bold=True, color=col)
        txt(slide, Inches(9.0),  y+Inches(0.07), Inches(3.8), Inches(0.54),
            desc, size=12, color=BLACK, wrap=True)

    txt(slide, Inches(0.4), Inches(6.58), Inches(12.6), Inches(0.28),
        "Cost: just 15% random dropout on 4 dims per batch. No extra parameters. No second model.",
        size=11, italic=True, color=MGRAY)


def s3_arch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    header(slide, "U-Net Architecture — The Denoiser Network",
           "Encoder compresses, decoder recovers, skip connections preserve detail")
    pill(slide, Inches(11.9), Inches(0.12), "4 / 7", GREEN)

    txt(slide, Inches(0.4), Inches(1.12), Inches(5.6), Inches(0.3),
        "Inputs", size=14, bold=True, color=NAVY)
    io_rows = [
        (LBLUE,  NAVY,   "Noisy action  at",        "(B, 32 steps, 5 dims)"),
        (LGREEN, GREEN,  "Observation  obs",          "(B, 26 dims)"),
        (LPURP,  PURPLE, "Timestep  t",               "(B,)  how noisy?"),
        (LGOLD,  GOLD,   "Output: noise prediction",  "(B, 32, 5)  same as at"),
    ]
    for i, (fill, col, name, shape) in enumerate(io_rows):
        y = Inches(1.44) + i*Inches(0.6)
        rect(slide, Inches(0.4), y, Inches(5.6), Inches(0.5), fill, col, Pt(1))
        txt(slide, Inches(0.55), y+Inches(0.06), Inches(2.3), Inches(0.38),
            name, size=12, bold=True, color=col)
        txt(slide, Inches(2.9), y+Inches(0.06), Inches(2.9), Inches(0.38),
            shape, size=11, color=BLACK)

    fcard(slide, Inches(0.4), Inches(3.95), Inches(5.6), Inches(0.68),
          "Time Embedding   (encodes noise level t as a vector)",
          "e_t = [sin(t/freq_k), cos(t/freq_k)]  then MLP -> 256 dims")

    fcard(slide, Inches(0.4), Inches(4.73), Inches(5.6), Inches(0.65),
          "Input Fusion   (merge action and observation)",
          "x  =  Linear(at)  +  MLP(obs)      both -> 256 dims")

    fcard(slide, Inches(0.4), Inches(5.48), Inches(5.6), Inches(0.68),
          "Each Block   (linear + norm + activation + skip)",
          "h = Mish( Norm(Linear(x)) + time_proj(e_t) )"
          "\nout = Mish( Norm(Linear(h)) + shortcut(x) )")

    txt(slide, Inches(6.4), Inches(1.12), Inches(6.6), Inches(0.3),
        "Network Flow", size=14, bold=True, color=NAVY)
    blocks = [
        (TEAL,   "Input  x = Linear(at) + MLP(obs)",             "256 ch"),
        (BLUE,   "Encoder Block 1  +  time",                      "256->512"),
        (BLUE,   "Encoder Block 2  +  time   -> save skip",        "512->1024"),
        (NAVY,   "Bottleneck  (global context)",                   "1024->1024"),
        (GREEN,  "Decoder Block 1  +  skip from Enc2",            "2048->512"),
        (GREEN,  "Decoder Block 2  +  skip from Enc1",            "1024->256"),
        (ORANGE, "Output Linear  ->  predicted noise",             "256->5"),
    ]
    for i, (col, label, dim) in enumerate(blocks):
        y = Inches(1.46) + i*Inches(0.72)
        rect(slide, Inches(6.4), y, Inches(6.6), Inches(0.62), col)
        txt(slide, Inches(6.52), y+Inches(0.05), Inches(4.5), Inches(0.52),
            label, size=11, color=WHITE)
        txt(slide, Inches(11.2), y+Inches(0.14), Inches(1.7), Inches(0.34),
            dim, size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
        if i < len(blocks)-1:
            rect(slide, Inches(9.3), y+Inches(0.62), Inches(0.06), Inches(0.1), MGRAY)

    txt(slide, Inches(0.4), Inches(6.58), Inches(12.6), Inches(0.28),
        "Skip connections: encoder detail goes directly to the decoder -> sharp noise predictions.",
        size=11, italic=True, color=MGRAY)


def s4_ddim(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    header(slide, "DDIM Sampling — Fast Action Generation at Test Time",
           "20 steps instead of 100  |  eta=0.5 for diversity  |  K=4 candidates")
    pill(slide, Inches(11.9), Inches(0.12), "5 / 7", ORANGE)

    txt(slide, Inches(0.4), Inches(1.12), Inches(5.8), Inches(0.3),
        "Why DDIM?", size=14, bold=True, color=NAVY)
    txt(slide, Inches(0.4), Inches(1.44), Inches(5.8), Inches(0.5),
        "DDPM trained 100 steps but we only NEED 20 steps at inference. "
        "DDIM skips levels with negligible quality loss -> 5x faster.",
        size=13, color=BLACK, wrap=True)

    fcard(slide, Inches(0.4), Inches(2.1), Inches(5.8), Inches(0.72),
          "Step 1 — Estimate clean action from current noisy version",
          "a0_hat  =  ( at  -  sqrt(1 - alpha_t) * e_hat )  /  sqrt(alpha_t)")

    fcard(slide, Inches(0.4), Inches(2.94), Inches(5.8), Inches(0.82),
          "Step 2 — Move to the previous (less noisy) level",
          "a_{t-1}  =  sqrt(alpha_{t-1}) * a0_hat"
          "  +  sqrt(1 - alpha_{t-1} - sigma^2) * e_hat  +  sigma * z"
          "\nsigma = eta * sqrt(...)     z ~ N(0,I)     eta=0.5")

    fcard(slide, Inches(0.4), Inches(3.88), Inches(5.8), Inches(0.65),
          "CFG applied at EVERY step (steers toward behavior)",
          "e_hat  =  e_uncond  +  lambda * ( e_cond - e_uncond )")

    fcard(slide, Inches(0.4), Inches(4.65), Inches(5.8), Inches(0.6),
          "Denormalize final output back to real joint space",
          "a_real  =  a0_hat * std_actions  +  mean_actions")

    txt(slide, Inches(6.5), Inches(1.12), Inches(6.5), Inches(0.3),
        "Why K = 4 Candidates?", size=14, bold=True, color=NAVY)
    txt(slide, Inches(6.5), Inches(1.44), Inches(6.5), Inches(0.48),
        "eta=0.5 adds different random noise each run -> 4 "
        "different trajectories from the same observation.",
        size=13, color=BLACK, wrap=True)

    cands = [BLUE, GREEN, ORANGE, PURPLE]
    for i in range(4):
        y = Inches(2.05) + i*Inches(0.9)
        rect(slide, Inches(6.5), y, Inches(6.5), Inches(0.8), cands[i])
        txt(slide, Inches(6.62), y+Inches(0.12), Inches(6.2), Inches(0.56),
            f"Candidate {i+1}:  start pure noise  ->  20 DDIM steps  ->  action chunk (32 steps)",
            size=12, color=WHITE, wrap=True)

    rect(slide, Inches(6.5), Inches(5.65), Inches(6.5), Inches(0.6), GOLD)
    txt(slide, Inches(6.62), Inches(5.71), Inches(6.2), Inches(0.48),
        "All 4 sent to VLM scorer  ->  pick the highest-scoring one",
        size=13, bold=True, color=NAVY)

    txt(slide, Inches(0.4), Inches(6.58), Inches(12.6), Inches(0.28),
        "eta=0: same trajectory every time (no diversity).   "
        "eta=0.5: balanced variety for VLM selection.",
        size=11, italic=True, color=MGRAY)


def s5_vlm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    header(slide, "VLM Judging — Gemini 3 Pro Picks the Best Trajectory",
           "6 frames per candidate  ->  behavior question  ->  score 0 to 1")
    pill(slide, Inches(11.9), Inches(0.12), "6 / 7", RED)

    txt(slide, Inches(0.4), Inches(1.12), Inches(5.7), Inches(0.3),
        "Pipeline  (per episode)", size=14, bold=True, color=NAVY)
    pipe = [
        ("1", "Reset sim + place scene objects",  "Obstacle or waypoint per behavior type"),
        ("2", "Generate K=4 candidates",           "CFG-DDIM, eta=0.5"),
        ("3", "Capture 6 frames each candidate",  "t=0s 1s 2s 3s 4s 5s  (JPEG 240x240)"),
        ("4", "Ask Gemini 3 Pro",                  "Frames + behavior question in one call"),
        ("5", "Parse score from JSON response",    'score: 0.91,  cue: "early left lean"'),
        ("6", "Run best + worst candidates",       "Record video, measure path metrics"),
    ]
    for i, (n, t_, note) in enumerate(pipe):
        stepbox(slide, Inches(0.4), Inches(1.45)+i*Inches(0.83), Inches(5.7), Inches(0.73), n, t_, note)

    txt(slide, Inches(6.5), Inches(1.12), Inches(6.5), Inches(0.3),
        "What Gemini Is Asked Per Behavior", size=14, bold=True, color=NAVY)

    bhs = [
        (LBLUE,  NAVY,   "Legibility",
         "How EARLY can you tell which block the arm is going for?\n"
         "0 = cannot tell at all    1 = obvious from frame 1"),
        (LGREEN, GREEN,  "Predictability",
         "How STRAIGHT and DIRECT is the path?\n"
         "0 = curved or erratic    1 = perfectly straight"),
        (LORAN,  ORANGE, "Safety",
         "How much SPACE does the arm keep from the obstacle?\n"
         "0 = nearly hits it    1 = wide arc around it"),
        (LPURP,  PURPLE, "Grounding",
         "Does the arm PASS NEAR the waypoint before picking?\n"
         "0 = ignores waypoint    1 = clearly visits it"),
    ]
    for i, (fill, col, name, q) in enumerate(bhs):
        y = Inches(1.45) + i*Inches(1.28)
        rect(slide, Inches(6.5), y, Inches(6.5), Inches(1.18), fill, col, Pt(1.5))
        txt(slide, Inches(6.63), y+Inches(0.07), Inches(6.2), Inches(0.32),
            name, size=14, bold=True, color=col)
        txt(slide, Inches(6.63), y+Inches(0.38), Inches(6.2), Inches(0.72),
            q, size=12, color=BLACK, wrap=True)

    txt(slide, Inches(0.4), Inches(6.58), Inches(12.6), Inches(0.28),
        "No Gemini fine-tuning. No reward function. VLM judges behavior from video like a human observer.",
        size=11, italic=True, color=MGRAY)


def s6_metrics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    header(slide, "Evaluation Metrics  —  Four Geometric Numbers",
           "Computed on the actual executed trajectory after VLM selection")
    pill(slide, Inches(11.9), Inches(0.12), "7 / 7", GOLD)

    data = [
        (LBLUE,  NAVY,   "L_early  (Legibility)",
         "L_early  =  avg P(correct goal | position)   over first 30% of path",
         "Uses distance to each goal block as Bayesian signal",
         "Higher = arm reveals goal earlier"),
        (LGREEN, GREEN,  "Path Efficiency  (Predictability)",
         "path_eff  =  straight-line dist / actual path length",
         "1.0 = perfectly straight,  below 1 = curved or backtracking",
         "Higher = more direct"),
        (LORAN,  ORANGE, "Clearance  (Safety)",
         "clearance  =  minimum distance to obstacle  over all steps",
         "Obstacle radius = 3.5 cm.  Clearance > 5 cm = safe",
         "Higher = safer"),
        (LPURP,  PURPLE, "Hover Dist  (Grounding)",
         "hover_dist  =  minimum XY distance to waypoint  over all steps",
         "Did the arm pass near the instructed landmark?",
         "Lower = better"),
    ]

    for i, (fill, col, name, formula, explain, verdict) in enumerate(data):
        cx = Inches(0.3) + (i%2)*Inches(6.55)
        cy = Inches(1.12) + (i//2)*Inches(2.75)
        w, h = Inches(6.2), Inches(2.6)

        rect(slide, cx, cy, w, h, fill, col, Pt(1.5))

        txt(slide, cx+Inches(0.14), cy+Inches(0.1), w-Inches(0.25), Inches(0.34),
            name, size=14, bold=True, color=col)

        rect(slide, cx+Inches(0.14), cy+Inches(0.5), w-Inches(0.28), Inches(0.58), WHITE)
        txt(slide, cx+Inches(0.2), cy+Inches(0.54), w-Inches(0.38), Inches(0.5),
            formula, size=12, color=BLACK, wrap=True)

        txt(slide, cx+Inches(0.14), cy+Inches(1.14), w-Inches(0.25), Inches(0.75),
            explain, size=12, color=BLACK, wrap=True)

        txt(slide, cx+Inches(0.14), cy+Inches(2.15), w-Inches(0.25), Inches(0.3),
            "-> " + verdict, size=12, bold=True, color=col)

    txt(slide, Inches(0.3), Inches(6.6), Inches(12.8), Inches(0.28),
        "Compare VLM-steered (best of K) vs baseline (worst of K) -- same model, same checkpoint, only selection differs.",
        size=11, italic=True, color=MGRAY)


def s7_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)

    txt(slide, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.62),
        "Summary", size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rect(slide, Inches(0.3), Inches(0.94), Inches(12.7), Inches(0.42), BLUE)
    for j, (hdr, x, w) in enumerate([
            ("Component",     Inches(0.4),  Inches(1.8)),
            ("What It Does",  Inches(2.25), Inches(3.9)),
            ("Key Formula",   Inches(6.2),  Inches(6.6))]):
        txt(slide, x+Inches(0.08), Inches(0.96), w, Inches(0.38),
            hdr, size=12, bold=True, color=WHITE)

    rows = [
        (TEAL,   "DDPM Training",    "Learn to predict noise added to demos",
                 "L = E[ || e - e_theta(at, t, obs) ||^2 ]"),
        (PURPLE, "CFG Dropout",      "15% hidden -> one model, any behavior",
                 "e_hat = e_uncond + lambda * (e_cond - e_uncond)"),
        (GREEN,  "U-Net",            "6-block encoder-decoder, skip connections",
                 "x = Linear(at) + MLP(obs),  256->512->1024->512->256->5"),
        (ORANGE, "DDIM Sampling",    "20 steps, eta=0.5, K=4 diverse candidates",
                 "a0_hat = (at - sqrt(1-alpha_t)*e_hat) / sqrt(alpha_t)"),
        (RED,    "VLM Judging",      "Gemini 3 Pro scores 6 frames per candidate",
                 "score in [0,1] per behavior  ->  select best"),
        (GOLD,   "Metrics",          "L_early, efficiency, clearance, hover dist",
                 "VLM-steered vs baseline (same model)"),
    ]
    for i, (col, comp, idea, formula) in enumerate(rows):
        y = Inches(1.38) + i*Inches(0.77)
        row_bg = RGBColor(0x1F, 0x35, 0x5C) if i%2==0 else RGBColor(0x16, 0x27, 0x47)
        rect(slide, Inches(0.3), y, Inches(12.7), Inches(0.68), row_bg)
        rect(slide, Inches(0.3), y, Inches(0.07), Inches(0.68), col)
        for val, x, w in [
                (comp,    Inches(2.25), Inches(3.9)),
                (formula, Inches(6.2),  Inches(6.6))]:
            txt(slide, x+Inches(0.08), y+Inches(0.13), w-Inches(0.1), Inches(0.46),
                val, size=11, color=WHITE, wrap=True)
        txt(slide, Inches(0.42), y+Inches(0.13), Inches(1.75), Inches(0.46),
            comp, size=12, bold=True, color=WHITE)

    txt(slide, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.35),
        "One model  .  One training run  .  Behavior selected at inference via lambda and obs conditioning  .  VLM judges without fine-tuning",
        size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for i, (fn, name) in enumerate([
            (s0_title,   "Title / Overview"),
            (s1_ddpm,    "DDPM Training"),
            (s2_cfg,     "CFG Trick"),
            (s3_arch,    "U-Net Architecture"),
            (s4_ddim,    "DDIM Inference"),
            (s5_vlm,     "VLM Judging"),
            (s6_metrics, "Metrics"),
            (s7_summary, "Summary"),
    ], 1):
        fn(prs)
        print(f"  [{i}/8] {name}")
    prs.save(str(PPTX_PATH))
    print(f"\nSaved -> {PPTX_PATH}")

if __name__ == "__main__":
    main()
